"""
PayGuard Pitch Deck Generator
Creates a 25-slide PowerPoint presentation using python-pptx.
Black and white color scheme, Calibri fonts, widescreen 16:9.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Emu
import copy
from lxml import etree

# ---------------------------------------------------------------------------
# Color constants
# ---------------------------------------------------------------------------
BLACK      = RGBColor(0,   0,   0)
WHITE      = RGBColor(255, 255, 255)
DARK_GRAY  = RGBColor(50,  50,  50)
MID_GRAY   = RGBColor(130, 130, 130)
LIGHT_GRAY = RGBColor(220, 220, 220)
PALE_GRAY  = RGBColor(240, 240, 240)
ALT_ROW    = RGBColor(245, 245, 245)

# ---------------------------------------------------------------------------
# Slide dimensions — widescreen 16:9
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.6)
CONTENT_LEFT  = MARGIN
CONTENT_TOP   = Inches(1.4)
CONTENT_W     = SLIDE_W - 2 * MARGIN
CONTENT_H     = SLIDE_H - CONTENT_TOP - Inches(0.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    """Return the blank slide layout (index 6)."""
    return prs.slide_layouts[6]


def set_slide_background(slide, r=255, g=255, b=255):
    """Fill slide background with a solid color (default white)."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_text_box(slide, text, left, top, width, height,
                 font_size=16, bold=False,
                 alignment=PP_ALIGN.LEFT,
                 font_color=None,
                 wrap=True,
                 italic=False,
                 font_name="Calibri"):
    """Add a text box and return the shape."""
    if font_color is None:
        font_color = BLACK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = font_color
    return txBox


def add_multiline_text_box(slide, lines, left, top, width, height,
                            font_size=16, bold=False,
                            alignment=PP_ALIGN.LEFT,
                            font_color=None,
                            line_space_pt=None,
                            font_name="Calibri"):
    """Add a text box with multiple paragraphs."""
    if font_color is None:
        font_color = BLACK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.name  = font_name
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = font_color
    return txBox


def add_stat_box(slide, stat, label, left, top, width, height,
                 stat_size=36, label_size=13):
    """Add a gray box with a large stat number and a label below it."""
    # Background rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PALE_GRAY
    line = shape.line
    line.color.rgb = BLACK
    line.width = Pt(1)

    # Stat number
    inner_pad = Inches(0.1)
    stat_h = height * 0.58
    stat_box = slide.shapes.add_textbox(
        left + inner_pad, top + Inches(0.05),
        width - 2 * inner_pad, stat_h
    )
    tf = stat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = stat
    run.font.name = "Calibri"
    run.font.size = Pt(stat_size)
    run.font.bold = True
    run.font.color.rgb = BLACK

    # Label
    lbl_box = slide.shapes.add_textbox(
        left + inner_pad, top + stat_h,
        width - 2 * inner_pad, height - stat_h
    )
    tf2 = lbl_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.name = "Calibri"
    run2.font.size = Pt(label_size)
    run2.font.bold = False
    run2.font.color.rgb = DARK_GRAY


def add_feature_box(slide, title, body_lines, left, top, width, height):
    """Add a bordered box with a bold title and bullet body lines."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PALE_GRAY
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.75)

    pad = Inches(0.15)
    # Title
    t_box = slide.shapes.add_textbox(
        left + pad, top + pad, width - 2 * pad, Inches(0.4)
    )
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BLACK

    # Body
    b_box = slide.shapes.add_textbox(
        left + pad, top + Inches(0.45), width - 2 * pad, height - Inches(0.6)
    )
    tf2 = b_box.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        run2 = p2.add_run()
        run2.text = line
        run2.font.name = "Calibri"
        run2.font.size = Pt(12)
        run2.font.color.rgb = DARK_GRAY


def add_table(slide, data, left, top, width, height,
              col_widths=None, header_font_size=12, body_font_size=11):
    """
    Add a styled table. data[0] is the header row.
    col_widths: list of fractions summing to 1.0 (or None for equal widths).
    """
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    if col_widths:
        for ci, frac in enumerate(col_widths):
            tbl.columns[ci].width = int(width * frac)

    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)

            # Styling
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"

            if ri == 0:
                # Header
                run.font.size  = Pt(header_font_size)
                run.font.bold  = True
                run.font.color.rgb = WHITE
                # black fill
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = BLACK
            else:
                run.font.size  = Pt(body_font_size)
                run.font.bold  = False
                run.font.color.rgb = BLACK
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = WHITE if ri % 2 == 1 else ALT_ROW

            # Borders
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                ln = etree.SubElement(tcPr, qn(border_tag.replace("a:", "a:")))
                ln.set("w", "6350")   # 0.5pt in EMUs
                ln.set("cap", "flat")
                ln.set("cmpd", "sng")
                solidFill = etree.SubElement(ln, qn("a:solidFill"))
                srgbClr  = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", "000000")

    return tbl


def add_slide_number(slide, num, total=25):
    """Add 'X / 25' to the bottom-right corner."""
    left  = SLIDE_W - Inches(1.3)
    top   = SLIDE_H - Inches(0.4)
    width = Inches(1.1)
    height = Inches(0.3)
    add_text_box(slide, f"{num} / {total}", left, top, width, height,
                 font_size=9, bold=False,
                 alignment=PP_ALIGN.RIGHT,
                 font_color=MID_GRAY)


def add_watermark(slide):
    """Add 'PayGuard' watermark text at bottom-left."""
    add_text_box(slide, "PayGuard",
                 MARGIN, SLIDE_H - Inches(0.4),
                 Inches(2.5), Inches(0.3),
                 font_size=9, bold=False,
                 alignment=PP_ALIGN.LEFT,
                 font_color=LIGHT_GRAY)


def add_slide_title(slide, title_text, font_size=28):
    """Add the bold slide title at the top."""
    title_box = add_text_box(
        slide, title_text,
        MARGIN, Inches(0.3),
        CONTENT_W, Inches(0.9),
        font_size=font_size, bold=True,
        alignment=PP_ALIGN.LEFT,
        font_color=BLACK
    )
    # Thin horizontal rule below title
    rule = slide.shapes.add_shape(
        1,
        MARGIN, Inches(1.1),
        CONTENT_W, Pt(1.5)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLACK
    rule.line.fill.background()
    return title_box


def add_divider_line(slide, top_pos):
    rule = slide.shapes.add_shape(1, MARGIN, top_pos, CONTENT_W, Pt(1))
    rule.fill.solid()
    rule.fill.fore_color.rgb = LIGHT_GRAY
    rule.line.fill.background()


# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------

def build_deck():
    prs = new_presentation()
    blank = blank_layout(prs)

    # -----------------------------------------------------------------------
    # Slide 1 — Title Slide
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)

    # Large "PayGuard" title
    add_text_box(slide, "PayGuard",
                 MARGIN, Inches(1.5), CONTENT_W, Inches(1.5),
                 font_size=64, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    # Horizontal rule
    rule = slide.shapes.add_shape(1, Inches(2), Inches(3.0), Inches(9.33), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLACK
    rule.line.fill.background()

    # Subtitle
    add_text_box(slide,
                 "AI-Powered ATM Monitoring & Customer Transparency Platform",
                 MARGIN, Inches(3.15), CONTENT_W, Inches(0.7),
                 font_size=22, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=DARK_GRAY)

    # Tag line
    add_text_box(slide,
                 "Transforming India's ATM Infrastructure",
                 MARGIN, Inches(3.9), CONTENT_W, Inches(0.55),
                 font_size=16, bold=False, italic=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # Company / date
    add_text_box(slide,
                 "PayGuard  |  March 2026",
                 MARGIN, Inches(5.4), CONTENT_W, Inches(0.4),
                 font_size=13, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # Bottom watermark
    add_watermark(slide)

    # -----------------------------------------------------------------------
    # Slide 2 — The Problem: Scale
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "India's ATM Crisis by the Numbers")
    add_watermark(slide)
    add_slide_number(slide, 2)

    # 4 stat boxes in a row
    box_w = Inches(2.8)
    box_h = Inches(1.9)
    box_top = Inches(1.55)
    gap = Inches(0.22)
    start_x = MARGIN

    stats = [
        ("255,000", "ATMs across India"),
        ("1.05 Billion", "Transactions / month"),
        ("5–8%", "ATM failure rate"),
        ("52.5M+", "Failed transactions / month"),
    ]
    for i, (stat, lbl) in enumerate(stats):
        add_stat_box(slide, stat, lbl,
                     start_x + i * (box_w + gap), box_top,
                     box_w, box_h, stat_size=30, label_size=12)

    # Bottom line
    add_text_box(slide,
                 "That's 1.75 million failed transactions every single day",
                 MARGIN, Inches(3.7), CONTENT_W, Inches(0.5),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_multiline_text_box(slide, [
        "ATM uptime in India: 90–93%  vs  98%+ in developed markets",
        "70% of failures go undetected for more than 30 minutes",
        "Each failed ATM costs ₹15,000–50,000 per day in lost revenue & penalties"
    ],
        MARGIN, Inches(4.3), CONTENT_W, Inches(1.8),
        font_size=14, alignment=PP_ALIGN.LEFT, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 3 — Customer Impact
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "What Happens When an ATM Fails?")
    add_watermark(slide)
    add_slide_number(slide, 3)

    pain_points = [
        ("No Transparency",
         ["Customer has no idea if ATM is", "down, network issue, or hardware", "failure."]),
        ("Overwhelmed Helplines",
         ["Call centres flooded with ATM", "failure complaints; average wait:", ">20 minutes."]),
        ("No Refund Tracking",
         ["Customers cannot track refund", "progress or expected resolution", "timeline."]),
        ("Zero Engineer Visibility",
         ["Banks cannot see which ATMs", "need immediate dispatch vs.", "remote fix."]),
    ]

    box_w = Inches(2.85)
    box_h = Inches(2.5)
    gap   = Inches(0.2)
    top   = Inches(1.55)
    for i, (title, lines) in enumerate(pain_points):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (box_w + gap), top,
                        box_w, box_h)

    # Stat highlight
    add_text_box(slide,
                 "73% of Indian ATM users have experienced a failed transaction  (NPCI 2023)",
                 MARGIN, Inches(4.25), CONTENT_W, Inches(0.5),
                 font_size=15, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_text_box(slide,
                 "Current resolution model: customer notices → calls helpline → ticket raised → engineer dispatched manually → 4+ hour MTTR",
                 MARGIN, Inches(4.9), CONTENT_W, Inches(0.7),
                 font_size=13, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 4 — Operator Impact
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "The Operator Blind Spot")
    add_watermark(slide)
    add_slide_number(slide, 4)

    stats4 = [
        ("₹15K–50K", "Downtime cost per ATM/day"),
        ("70%", "Failures undetected >30 min"),
        ("₹630 Cr", "ATM fraud FY 2022-23"),
        ("4+ hrs", "Average MTTR today"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.9)
    for i, (stat, lbl) in enumerate(stats4):
        add_stat_box(slide, stat, lbl,
                     MARGIN + i * (box_w + gap), Inches(1.55),
                     box_w, box_h, stat_size=28, label_size=12)

    add_text_box(slide,
                 "Only 30% of failures are detected within 30 minutes",
                 MARGIN, Inches(3.7), CONTENT_W, Inches(0.45),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    bullets = [
        "Banks rely on customer complaints to discover ATM failures — fully reactive",
        "Legacy SNMP monitoring tools lack AI, predictive capability, or customer-facing layer",
        "No automated self-healing: every incident requires manual engineer dispatch",
        "Incident management spreadsheets and siloed ticketing systems create blind spots",
    ]
    add_multiline_text_box(slide, bullets,
                           MARGIN, Inches(4.3), CONTENT_W, Inches(2.5),
                           font_size=13, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 5 — Regulatory Pressure
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "RBI Is Watching")
    add_watermark(slide)
    add_slide_number(slide, 5)

    regs = [
        ("RBI Circular RBI/2011-12/314",
         ["5-working-day refund mandate for failed ATM", "transactions — automatically triggered", "by RBI regulations."]),
        ("Penalty: ₹100/day/customer",
         ["For every day delay beyond the mandate,", "banks pay ₹100 per affected customer —", "untracked & unmanaged today."]),
        ("DPDP Act 2023",
         ["Digital Personal Data Protection Act", "mandates SHA-256 phone hashing and", "data minimisation for customer portals."]),
        ("PCI DSS v4.0",
         ["Payment Card Industry standards alignment", "for all ATM transaction processing", "and fraud detection workflows."]),
    ]
    box_w = Inches(2.85)
    box_h = Inches(2.5)
    for i, (title, lines) in enumerate(regs):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (box_w + gap), Inches(1.55),
                        box_w, box_h)

    add_text_box(slide,
                 "\"Banks lack automated compliance workflows — PayGuard builds RBI compliance in, not on.\"",
                 MARGIN, Inches(4.3), CONTENT_W, Inches(0.6),
                 font_size=14, bold=True, italic=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_text_box(slide,
                 "Regulatory pressure is our sales team",
                 MARGIN, Inches(5.0), CONTENT_W, Inches(0.4),
                 font_size=14, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 6 — Market Opportunity
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "A USD 18.2 Billion Market with No AI-First Player")
    add_watermark(slide)
    add_slide_number(slide, 6)

    stats6 = [
        ("$18.2B", "Global ATM Managed\nServices by 2028\n(8.3% CAGR)"),
        ("$220M", "India ATM Monitoring\nSoftware Addressable\nMarket"),
        ("₹32,000 Cr", "Indian Banking IT\nSpend / Year"),
        ("$11–18M", "3-Year Realistic\nARR Capture"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(2.2)
    for i, (stat, lbl) in enumerate(stats6):
        add_stat_box(slide, stat, lbl,
                     MARGIN + i * (box_w + gap), Inches(1.55),
                     box_w, box_h, stat_size=26, label_size=11)

    add_text_box(slide,
                 "No existing player combines AI monitoring + predictive failure + customer transparency portal",
                 MARGIN, Inches(4.0), CONTENT_W, Inches(0.45),
                 font_size=14, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_multiline_text_box(slide, [
        "India is the world's 2nd largest ATM network by count, yet monitoring software market is severely underserved",
        "Global players (NCR, Diebold) are not optimised for Indian banking regulations, languages, or infrastructure",
    ],
        MARGIN, Inches(4.6), CONTENT_W, Inches(1.4),
        font_size=13, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 7 — ATM Network Breakdown
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "India's ATM Landscape — The Addressable Fleet")
    add_watermark(slide)
    add_slide_number(slide, 7)

    table_data = [
        ["Bank Category",         "No. of ATMs",  "Key Players",                          "Priority"],
        ["Public Sector Banks",   "150,000",      "SBI, PNB, Bank of Baroda, Canara",     "High"],
        ["Private Sector Banks",  "75,000",       "HDFC, ICICI, Axis, Kotak",             "Very High"],
        ["White-Label Operators", "30,000",       "Tata Communications, Hitachi, AGS",    "Medium"],
        ["Total",                 "255,000",      "12 PSBs + 22 Private + WLOs",          "—"],
    ]
    add_table(slide, table_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(2.5),
              col_widths=[0.26, 0.18, 0.38, 0.18])

    add_text_box(slide,
                 "12 public sector banks.  22 private sector banks.  One platform to monitor them all.",
                 MARGIN, Inches(4.2), CONTENT_W, Inches(0.45),
                 font_size=15, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_multiline_text_box(slide, [
        "Go-to-market priority: Private sector banks first (faster procurement cycles, tech-forward culture)",
        "PSB expansion in Phase 2 via system integrator partnerships (TCS, Infosys, Wipro)",
        "White-label operators: secondary channel through revenue-share agreements"
    ],
        MARGIN, Inches(4.8), CONTENT_W, Inches(1.6),
        font_size=13, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 8 — Solution Overview
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Introducing PayGuard")
    add_watermark(slide)
    add_slide_number(slide, 8)

    # Two columns
    col_w = Inches(5.6)
    col_h = Inches(3.5)
    col_top = Inches(1.55)

    # PayGuard box
    pulse_box = slide.shapes.add_shape(1, MARGIN, col_top, col_w, col_h)
    pulse_box.fill.solid()
    pulse_box.fill.fore_color.rgb = PALE_GRAY
    pulse_box.line.color.rgb = BLACK
    pulse_box.line.width = Pt(1.5)

    add_text_box(slide, "PayGuard — For Banks & Operators",
                 MARGIN + Inches(0.1), col_top + Inches(0.1), col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, font_color=BLACK)

    pulse_features = [
        "Real-time ATM fleet health monitoring",
        "AI root cause classifier (8 categories)",
        "Predictive failure detection",
        "Anomaly & fraud detection (Z-score)",
        "Self-healing automation",
        "Incident management workflow",
        "Engineer dispatch dashboard",
    ]
    add_multiline_text_box(slide, ["• " + f for f in pulse_features],
                           MARGIN + Inches(0.15), col_top + Inches(0.65),
                           col_w - Inches(0.3), col_h - Inches(0.8),
                           font_size=12, font_color=DARK_GRAY)

    # PayGuard box
    pg_left = MARGIN + col_w + Inches(0.53)
    pg_box = slide.shapes.add_shape(1, pg_left, col_top, col_w, col_h)
    pg_box.fill.solid()
    pg_box.fill.fore_color.rgb = WHITE
    pg_box.line.color.rgb = BLACK
    pg_box.line.width = Pt(1.5)

    add_text_box(slide, "PayGuard — For Customers",
                 pg_left + Inches(0.1), col_top + Inches(0.1), col_w - Inches(0.2), Inches(0.5),
                 font_size=15, bold=True, font_color=BLACK)

    pg_features = [
        "OTP-authenticated customer portal",
        "Real-time incident status tracking",
        "Refund progress bar with ETA",
        "Bilingual: English + Hindi",
        "WebSocket live push updates",
        "Automated SMS notifications",
        "DPDP Act 2023 compliant",
    ]
    add_multiline_text_box(slide, ["• " + f for f in pg_features],
                           pg_left + Inches(0.15), col_top + Inches(0.65),
                           col_w - Inches(0.3), col_h - Inches(0.8),
                           font_size=12, font_color=DARK_GRAY)

    # Architecture flow
    add_text_box(slide,
                 "ATM Fleet  →  PayGuard Engine (Django)  →  FastAPI AI Service  →  Engineer Dashboard  /  Customer Portal",
                 MARGIN, Inches(5.2), CONTENT_W, Inches(0.4),
                 font_size=12, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 9 — Real-Time Monitoring
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 1: Real-Time ATM Fleet Monitoring")
    add_watermark(slide)
    add_slide_number(slide, 9)

    dims = [
        ("Network Score",     ["Connectivity quality,", "latency, packet loss,", "uptime metrics"]),
        ("Hardware Score",    ["Card reader, cash unit,", "receipt printer,", "physical sensor data"]),
        ("Software Score",    ["OS health, service", "process status,", "config integrity"]),
        ("Transaction Score", ["Success rate, timeout", "rate, reversal rate,", "throughput"]),
    ]
    box_w = Inches(2.85)
    box_h = Inches(2.4)
    for i, (title, lines) in enumerate(dims):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (box_w + gap), Inches(1.55),
                        box_w, box_h)

    add_text_box(slide,
                 "Composite health score per ATM, updated in real-time via WebSocket (Django Channels)",
                 MARGIN, Inches(4.15), CONTENT_W, Inches(0.45),
                 font_size=14, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    status_data = [
        ["Status",    "Meaning",                           "Automated Action"],
        ["ONLINE",    "Health score ≥ 80",                 "Monitor continuously"],
        ["DEGRADED",  "Health score 50–79",                "Raise alert, predictive scan"],
        ["OFFLINE",   "Health score < 50 / unreachable",   "Trigger self-heal + incident"],
        ["MAINTENANCE","Scheduled downtime",               "Suppress alerts, log event"],
    ]
    add_table(slide, status_data,
              MARGIN, Inches(4.75), CONTENT_W, Inches(1.95),
              col_widths=[0.18, 0.37, 0.45])

    # -----------------------------------------------------------------------
    # Slide 10 — AI Root Cause Classifier
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 2: AI Root Cause Classifier")
    add_watermark(slide)
    add_slide_number(slide, 10)

    add_text_box(slide,
                 "Input: ATM error log event codes  →  Output: Root cause category + Confidence score + Recommended self-heal action",
                 MARGIN, Inches(1.3), CONTENT_W, Inches(0.4),
                 font_size=12, bold=False, italic=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    cat_data = [
        ["Category",   "Description",                              "Typical Self-Heal"],
        ["NETWORK",    "Connectivity failures, latency spikes",    "SWITCH_NETWORK"],
        ["CASH_JAM",   "Cassette jam, dispenser error",           "ALERT_ENGINEER"],
        ["HARDWARE",   "Card reader, printer, physical fault",     "ALERT_ENGINEER"],
        ["SERVER",     "Core banking timeout, host unreachable",   "RESTART_SERVICE"],
        ["FRAUD",      "Suspicious transaction pattern detected",  "FREEZE_ATM"],
        ["TIMEOUT",    "Network / session timeout",                "FLUSH_CACHE"],
        ["SWITCH",     "Switch / routing layer failure",           "RESTART_SERVICE"],
        ["UNKNOWN",    "Confidence < 0.6, unclassified",          "ALERT_ENGINEER"],
    ]
    add_table(slide, cat_data,
              MARGIN, Inches(1.85), CONTENT_W, Inches(3.5),
              col_widths=[0.2, 0.5, 0.3])

    add_text_box(slide,
                 "Confidence scoring on every classification — no black boxes.  Threshold: 0.6 (below = UNKNOWN)",
                 MARGIN, Inches(5.55), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER,
                 font_color=BLACK)

    add_text_box(slide,
                 "FastAPI microservice  |  Lookup-table + ML ensemble  |  <50ms response time",
                 MARGIN, Inches(6.05), CONTENT_W, Inches(0.35),
                 font_size=12, bold=False,
                 alignment=PP_ALIGN.CENTER,
                 font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 11 — Predictive Failure Detection
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 3: Predictive Failure Detection")
    add_watermark(slide)
    add_slide_number(slide, 11)

    add_text_box(slide, "From Reactive to Predictive",
                 MARGIN, Inches(1.45), CONTENT_W, Inches(0.45),
                 font_size=18, bold=True,
                 alignment=PP_ALIGN.LEFT, font_color=BLACK)

    add_multiline_text_box(slide, [
        "The PayGuard AI service ingests health snapshot time-series for every ATM and applies:",
        "",
        "  Rolling Mean: baseline health score over configurable window",
        "  Rolling Variance: volatility indicator — high variance precedes failure",
        "  Slope Analysis: linear regression on recent health trajectory",
        "  Failure Probability Score: composite 0.0–1.0 score per ATM",
        "",
        "When slope is negative AND variance is high, the system raises a PREDICTIVE alert",
        "before the ATM actually fails — giving engineers a head start.",
    ],
        MARGIN, Inches(2.0), Inches(8.5), Inches(3.8),
        font_size=13, font_color=DARK_GRAY)

    # Right panel
    r_left = Inches(9.3)
    r_top  = Inches(1.55)
    r_w    = Inches(3.6)
    r_h    = Inches(4.4)
    r_box  = slide.shapes.add_shape(1, r_left, r_top, r_w, r_h)
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = PALE_GRAY
    r_box.line.color.rgb = BLACK
    r_box.line.width = Pt(1)

    add_text_box(slide, "Key Metrics",
                 r_left + Inches(0.15), r_top + Inches(0.1), r_w - Inches(0.3), Inches(0.4),
                 font_size=13, bold=True, font_color=BLACK)

    key_metrics = [
        "Prediction horizon: 2–24 hours",
        "Time-series window: 48 snapshots",
        "Slope threshold: -0.5/hr",
        "Variance threshold: σ² > 25",
        "Alert lead time: avg 3.2 hours",
        "False positive rate: <8%",
    ]
    add_multiline_text_box(slide, ["• " + m for m in key_metrics],
                           r_left + Inches(0.15), r_top + Inches(0.6),
                           r_w - Inches(0.3), r_h - Inches(0.8),
                           font_size=12, font_color=DARK_GRAY)

    add_text_box(slide,
                 "Identifies declining ATMs before they reach zero — predicts failure window in hours, not after the fact",
                 MARGIN, Inches(6.1), CONTENT_W, Inches(0.5),
                 font_size=14, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 12 — Anomaly & Fraud Detection
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 4: AI Anomaly & Fraud Detection")
    add_watermark(slide)
    add_slide_number(slide, 12)

    # Z-score anomaly
    left_w = Inches(6.0)
    left_box = slide.shapes.add_shape(1, MARGIN, Inches(1.55), left_w, Inches(4.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = PALE_GRAY
    left_box.line.color.rgb = BLACK
    left_box.line.width = Pt(1)

    add_text_box(slide, "Statistical Anomaly Detection",
                 MARGIN + Inches(0.15), Inches(1.65), left_w - Inches(0.3), Inches(0.45),
                 font_size=14, bold=True, font_color=BLACK)

    add_multiline_text_box(slide, [
        "Z-score analysis on error rate vs. historical baseline:",
        "",
        "  Z = (observed_rate - mean) / std_dev",
        "",
        "  Threshold: σ > 3.0  →  ANOMALY_FLAG raised",
        "",
        "  Types detected:",
        "    • RAPID_FAILURES — burst error pattern",
        "    • MALWARE_PATTERN — unusual code sequences",
        "    • UNUSUAL_WITHDRAWAL — cash pattern outlier",
        "    • CARD_SKIMMING — read/write anomaly",
    ],
        MARGIN + Inches(0.15), Inches(2.15), left_w - Inches(0.3), Inches(3.8),
        font_size=12, font_color=DARK_GRAY)

    # Fraud heuristics
    r_left2 = MARGIN + left_w + Inches(0.4)
    r_w2    = CONTENT_W - left_w - Inches(0.4)
    r_box2  = slide.shapes.add_shape(1, r_left2, Inches(1.55), r_w2, Inches(4.5))
    r_box2.fill.solid()
    r_box2.fill.fore_color.rgb = WHITE
    r_box2.line.color.rgb = BLACK
    r_box2.line.width = Pt(1)

    add_text_box(slide, "Multi-Heuristic Fraud Detection",
                 r_left2 + Inches(0.15), Inches(1.65), r_w2 - Inches(0.3), Inches(0.45),
                 font_size=14, bold=True, font_color=BLACK)

    fraud_bullets = [
        "1. Rapid Withdrawal",
        "   3+ transactions in a 10-minute",
        "   window from same card/ATM",
        "",
        "2. Amount Anomaly",
        "   Z-score > 3.5σ from card's",
        "   own transaction baseline",
        "",
        "3. Geographic Anomaly",
        "   >100 km in <60 minutes",
        "   (impossible travel detection)",
        "",
        "Action: FREEZE_ATM + ALERT_ENGINEER",
    ]
    add_multiline_text_box(slide, fraud_bullets,
                           r_left2 + Inches(0.15), Inches(2.15),
                           r_w2 - Inches(0.3), Inches(3.8),
                           font_size=12, font_color=DARK_GRAY)

    add_text_box(slide,
                 "Only platform combining statistical anomaly detection with multi-heuristic fraud analysis",
                 MARGIN, Inches(6.3), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 13 — Self-Healing Automation
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 5: Self-Healing Automation")
    add_watermark(slide)
    add_slide_number(slide, 13)

    heal_data = [
        ["Root Cause",    "Automated Action",    "Description",                                    "Avg Resolution"],
        ["NETWORK",       "SWITCH_NETWORK",      "Failover to backup connectivity (4G/VSAT)",      "< 5 minutes"],
        ["SERVER/SWITCH", "RESTART_SERVICE",     "Remote service restart via secure API",           "< 10 minutes"],
        ["TIMEOUT",       "FLUSH_CACHE",         "Clear session/connection cache remotely",         "< 3 minutes"],
        ["FRAUD",         "FREEZE_ATM",          "Suspend card ops + notify fraud team",           "Immediate"],
        ["HARDWARE",      "ALERT_ENGINEER",      "Auto-dispatch nearest engineer with context",    "< 30 minutes"],
        ["CASH_JAM",      "ALERT_ENGINEER",      "Dispatch cash-handling technician",              "< 45 minutes"],
        ["UNKNOWN",       "ALERT_ENGINEER",      "Escalate with full log context bundle",          "< 30 minutes"],
    ]
    add_table(slide, heal_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(4.0),
              col_widths=[0.2, 0.22, 0.38, 0.2])

    add_text_box(slide,
                 "MTTR reduced from 4+ hours to under 30 minutes for software/network failures",
                 MARGIN, Inches(5.75), CONTENT_W, Inches(0.45),
                 font_size=15, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    add_text_box(slide,
                 "All self-heal actions logged to SelfHealAction model with timestamp, outcome, and audit trail",
                 MARGIN, Inches(6.25), CONTENT_W, Inches(0.4),
                 font_size=12, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 14 — PayGuard Customer Portal
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Feature 6: PayGuard Customer Portal")
    add_watermark(slide)
    add_slide_number(slide, 14)

    add_text_box(slide,
                 "The First Customer-Facing ATM Incident Transparency Portal in India",
                 MARGIN, Inches(1.3), CONTENT_W, Inches(0.4),
                 font_size=15, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # Status flow
    statuses = [
        "DETECTED",
        "INVESTIGATING",
        "ENGINEER\nEN ROUTE",
        "RESOLVING",
        "REFUND\nINITIATED",
        "RESOLVED",
    ]
    status_w = Inches(1.85)
    status_h = Inches(0.75)
    status_top = Inches(1.9)
    status_gap = Inches(0.15)
    for i, s in enumerate(statuses):
        sx = MARGIN + i * (status_w + status_gap)
        s_box = slide.shapes.add_shape(1, sx, status_top, status_w, status_h)
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = BLACK if i == len(statuses) - 1 else PALE_GRAY
        s_box.line.color.rgb = BLACK
        s_box.line.width = Pt(1)

        add_text_box(slide, s, sx + Inches(0.05), status_top + Inches(0.05),
                     status_w - Inches(0.1), status_h - Inches(0.1),
                     font_size=10, bold=True,
                     alignment=PP_ALIGN.CENTER,
                     font_color=WHITE if i == len(statuses) - 1 else BLACK)

    # Arrow label
    add_text_box(slide, "← Real-time progress visible to customer via PayGuard portal & SMS →",
                 MARGIN, Inches(2.75), CONTENT_W, Inches(0.35),
                 font_size=11, bold=False, italic=True,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    # Features
    portal_features = [
        ("OTP Authentication",     ["No password required", "OTP via SMS to registered", "mobile — familiar UPI flow"]),
        ("Bilingual EN/HI",        ["Full English + Hindi UI", "Hindi: ATM विफलता सूचना", "RBI-mandate messages"]),
        ("Refund Progress Bar",    ["Visual timeline with ETA", "Auto-updates via WebSocket", "5-day RBI countdown"]),
        ("Live WebSocket Updates", ["Django Channels push", "Exponential backoff retry", "Mobile-optimised"]),
    ]
    box_w = Inches(2.85)
    box_h = Inches(2.3)
    for i, (title, lines) in enumerate(portal_features):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (box_w + gap), Inches(3.2),
                        box_w, box_h)

    add_text_box(slide,
                 "SHA-256 phone hashing  |  DPDP Act 2023 compliant  |  No PII stored in plaintext",
                 MARGIN, Inches(5.65), CONTENT_W, Inches(0.35),
                 font_size=11, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 15 — Competitive Landscape
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Competitive Landscape")
    add_watermark(slide)
    add_slide_number(slide, 15)

    comp_data = [
        ["Feature",             "NCR Activate", "Diebold\nNixdorf", "FSS",  "AGS\nTransact", "PayGuard"],
        ["ATM Monitoring",      "Yes",           "Yes",              "Yes",  "Yes",            "Yes"],
        ["AI Root Cause",       "Partial",       "No",               "No",   "No",             "Yes"],
        ["Predictive Failure",  "No",            "No",               "No",   "No",             "Yes"],
        ["Customer Portal",     "No",            "No",               "No",   "No",             "Yes"],
        ["Fraud Detection",     "Partial",       "No",               "No",   "No",             "Yes"],
        ["Self-Healing",        "No",            "No",               "No",   "No",             "Yes"],
        ["India-First Design",  "No",            "No",               "Yes",  "Yes",            "Yes"],
        ["Bilingual EN/HI",     "No",            "No",               "No",   "No",             "Yes"],
        ["RBI Compliance",      "No",            "No",               "Partial","No",           "Yes"],
    ]
    add_table(slide, comp_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(4.7),
              col_widths=[0.28, 0.14, 0.14, 0.1, 0.14, 0.2],
              header_font_size=11, body_font_size=10)

    add_text_box(slide,
                 "PayGuard is the ONLY platform combining AI monitoring + predictive failure + customer transparency",
                 MARGIN, Inches(6.4), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 16 — India-First Design
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "India-First Design Principles")
    add_watermark(slide)
    add_slide_number(slide, 16)

    india_features = [
        ("Bilingual EN / HI",
         ["Full Hindi Unicode support",
          "ATM \u0935\u093f\u092b\u0932\u0924\u093e \u0938\u0942\u091a\u0928\u093e (ATM failure notice)",
          "\u0906\u092a\u0915\u093e \u0930\u093f\u092b\u0902\u0921 5 \u0926\u093f\u0928 \u092e\u0947\u0902 (Refund in 5 days)",
          "All customer-facing text bilingual"]),
        ("OTP Authentication",
         ["No passwords — OTP via SMS",
          "Familiar UPI-style auth flow",
          "Works on 2G/3G networks",
          "Timeout-safe implementation"]),
        ("Indian Number Formats",
         ["Lakh/Crore formatting",
          "\u20b9 currency symbol throughout",
          "Indian date formats (DD/MM/YYYY)",
          "IST timezone awareness"]),
        ("RBI Compliance Built-in",
         ["5-day refund mandate automated",
          "\u20b9100/day penalty tracker",
          "Audit trail for regulators",
          "Not bolted on — core design"]),
        ("DPDP Act 2023",
         ["SHA-256 phone hashing",
          "Data minimisation principles",
          "No PII in plaintext storage",
          "Right to erasure support"]),
        ("Mobile-First Resilience",
         ["WebSocket exponential backoff",
          "Graceful offline degradation",
          "Low-bandwidth SMS fallback",
          "Works on Indian mobile internet"]),
    ]

    box_w = Inches(3.85)
    box_h = Inches(2.25)
    gap2  = Inches(0.25)
    for i, (title, lines) in enumerate(india_features):
        col = i % 3
        row = i // 3
        lx = MARGIN + col * (box_w + gap2)
        ly = Inches(1.55) + row * (box_h + Inches(0.12))
        add_feature_box(slide, title, lines, lx, ly, box_w, box_h)

    # -----------------------------------------------------------------------
    # Slide 17 — Technical Architecture
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Technical Architecture")
    add_watermark(slide)
    add_slide_number(slide, 17)

    add_text_box(slide, "API-first  |  Microservice AI engine  |  Horizontally scalable  |  WebSocket real-time",
                 MARGIN, Inches(1.3), CONTENT_W, Inches(0.38),
                 font_size=12, bold=False, italic=True,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    layers = [
        ("Layer 1 — Presentation",
         "React 19 + TypeScript + Vite + Redux Toolkit",
         "Engineer Dashboard (port 3000)  |  PayGuard Customer Portal",
         Inches(1.6)),
        ("Layer 2 — Backend API",
         "Django 6 REST Framework + Django Channels (WebSocket)",
         "JWT Auth  |  REST API (port 8000)  |  ws/dashboard/  |  ws/logs/<atm_id>/",
         Inches(3.0)),
        ("Layer 3 — AI Microservice",
         "FastAPI + Python (port 8001)",
         "POST /classify  |  POST /predict  |  POST /detect  |  <50ms latency",
         Inches(4.4)),
    ]

    layer_w = CONTENT_W
    layer_h = Inches(1.1)
    for label, tech, desc, ly in layers:
        box = slide.shapes.add_shape(1, MARGIN, ly, layer_w, layer_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PALE_GRAY
        box.line.color.rgb = BLACK
        box.line.width = Pt(1)

        add_text_box(slide, label,
                     MARGIN + Inches(0.15), ly + Inches(0.05),
                     Inches(3.0), Inches(0.35),
                     font_size=12, bold=True, font_color=BLACK)
        add_text_box(slide, tech,
                     MARGIN + Inches(3.2), ly + Inches(0.05),
                     layer_w - Inches(3.4), Inches(0.35),
                     font_size=12, bold=False, font_color=DARK_GRAY)
        add_text_box(slide, desc,
                     MARGIN + Inches(0.15), ly + Inches(0.5),
                     layer_w - Inches(0.3), Inches(0.5),
                     font_size=11, bold=False, font_color=MID_GRAY)

        # Arrow down (except last)
        if ly < Inches(4.4):
            arr = slide.shapes.add_shape(1,
                MARGIN + layer_w / 2 - Inches(0.05),
                ly + layer_h,
                Inches(0.1), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = BLACK
            arr.line.fill.background()

    add_text_box(slide,
                 "Database: SQLite (dev) / PostgreSQL (prod)  |  Message broker: Redis (Channels layer)  |  Deployment: Docker + Gunicorn/Daphne",
                 MARGIN, Inches(5.75), CONTENT_W, Inches(0.45),
                 font_size=11, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    add_text_box(slide,
                 "Modern API-first architecture vs legacy SNMP polling used by all competitors",
                 MARGIN, Inches(6.25), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 18 — Business Model
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Business Model — SaaS Subscription Pricing")
    add_watermark(slide)
    add_slide_number(slide, 18)

    pricing_data = [
        ["Tier",           "Price",              "ATM Monitoring", "AI + Self-Heal", "Customer Portal", "SLA / Support"],
        ["Starter",        "₹800 / ATM / month", "Yes",            "No",             "No",              "Business hours"],
        ["Professional",   "₹1,800 / ATM / month","Yes",           "Yes",            "No",              "24/7 email"],
        ["Enterprise",     "₹3,500 / ATM / month","Yes",           "Yes",            "Yes",             "24/7 dedicated"],
    ]
    add_table(slide, pricing_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(2.2),
              col_widths=[0.17, 0.22, 0.16, 0.16, 0.16, 0.13])

    add_text_box(slide,
                 "A bank with 5,000 ATMs at Professional tier = ₹10.8 Crore / year ARR",
                 MARGIN, Inches(4.0), CONTENT_W, Inches(0.45),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    add_divider_line(slide, Inches(4.6))

    add_text_box(slide, "Additional Revenue Streams",
                 MARGIN, Inches(4.75), Inches(4.0), Inches(0.4),
                 font_size=14, bold=True, font_color=BLACK)

    streams = [
        "Professional Services: implementation, training, custom integration — ₹25-50 lakh/bank",
        "Data Insights API: anonymised fleet health data to ATM manufacturers and insurer — usage-based",
        "SMS Gateway: pass-through at margin for customer notifications",
        "Compliance Reports: automated RBI audit reports — ₹5 lakh/year add-on",
    ]
    add_multiline_text_box(slide, ["• " + s for s in streams],
                           MARGIN, Inches(5.25), CONTENT_W, Inches(1.8),
                           font_size=12, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 19 — Unit Economics
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Unit Economics")
    add_watermark(slide)
    add_slide_number(slide, 19)

    ue_stats = [
        ("₹50–80L", "Customer Acquisition\nCost (CAC)"),
        ("5–7 yrs", "Average Contract\nLength (LTV driver)"),
        ("70–75%", "Gross Margin\n(SaaS software)"),
        ("12–18 mo", "Payback Period\n(5,000 ATM bank)"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(2.0)
    for i, (stat, lbl) in enumerate(ue_stats):
        add_stat_box(slide, stat, lbl,
                     MARGIN + i * (box_w + gap), Inches(1.55),
                     box_w, box_h, stat_size=26, label_size=12)

    add_text_box(slide, "Banking software = long contracts, low churn",
                 MARGIN, Inches(3.75), CONTENT_W, Inches(0.4),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    ue_details = [
        ["Metric",                "Value",               "Notes"],
        ["ARR per 5K ATM bank",   "₹10.8 Cr",           "Professional tier"],
        ["ARR per 10K ATM bank",  "₹21.6 Cr",           "Professional tier"],
        ["Gross margin",          "70–75%",              "Infrastructure + support costs deducted"],
        ["NRR (net revenue ret.)", "110–120%",           "Upsell Starter→Pro→Enterprise"],
        ["Churn target",          "< 5% / year",        "Banking = 3-7 year vendor lock-in typical"],
        ["Break-even",            "Month 18–22",        "~12,000 ATMs under management"],
    ]
    add_table(slide, ue_details,
              MARGIN, Inches(4.3), CONTENT_W, Inches(2.7),
              col_widths=[0.35, 0.25, 0.40])

    # -----------------------------------------------------------------------
    # Slide 20 — Go-To-Market
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Go-To-Market Strategy")
    add_watermark(slide)
    add_slide_number(slide, 20)

    phases = [
        ("Phase 1: 0–12 Months",
         ["2 private sector bank pilots",
          "Geography: Chennai + Bengaluru",
          "Focus: 500–1,000 ATMs each",
          "Goal: prove ROI, reduce MTTR by 80%",
          "Channel: direct sales + demo",
          "Entry: CTO / VP-Technology office"]),
        ("Phase 2: 12–24 Months",
         ["Expand to PSBs (SBI, PNB, Canara)",
          "White-label ATM operators",
          "System integrator partnerships",
          "TCS, Infosys, Wipro channels",
          "Target: 15,000 ATMs total",
          "Launch PayGuard customer portal"]),
        ("Phase 3: 24–36 Months",
         ["Southeast Asia expansion",
          "Bangladesh: 12,000 ATMs",
          "Sri Lanka: 4,500 ATMs",
          "Nepal: 3,000 ATMs",
          "Localise for BDT, LKR, NPR",
          "Target: 50,000 ATMs total"]),
    ]

    col_w = Inches(3.9)
    col_h = Inches(3.8)
    col_top = Inches(1.55)
    col_gap = Inches(0.28)
    for i, (title, lines) in enumerate(phases):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (col_w + col_gap), col_top,
                        col_w, col_h)

    add_text_box(slide,
                 "Sales cycle: 6–12 months for private banks.  18–24 months for PSBs.  SI channel accelerates PSB access.",
                 MARGIN, Inches(5.55), CONTENT_W, Inches(0.4),
                 font_size=13, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    add_text_box(slide,
                 "Regulatory tailwind: RBI mandate enforcement creates urgency — compliance budget is non-discretionary",
                 MARGIN, Inches(6.0), CONTENT_W, Inches(0.45),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 21 — 5-Year Financial Projections
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "5-Year Financial Projections")
    add_watermark(slide)
    add_slide_number(slide, 21)

    proj_data = [
        ["Year",  "ATMs Managed", "Banks",  "ARR",        "Gross Margin"],
        ["Y1",    "2,000",        "2",      "₹4.3 Cr",   "65% (ramp)"],
        ["Y2",    "15,000",       "8",      "₹27 Cr",    "70%"],
        ["Y3",    "50,000",       "20",     "₹90 Cr",    "72%"],
        ["Y4",    "1,20,000",     "45",     "₹216 Cr",   "74%"],
        ["Y5",    "2,50,000",     "90",     "₹450 Cr",   "75%"],
    ]
    add_table(slide, proj_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(2.5),
              col_widths=[0.1, 0.2, 0.12, 0.2, 0.38])

    add_text_box(slide,
                 "Break-even at Month 18–22  (~12,000 ATMs under management)",
                 MARGIN, Inches(4.2), CONTENT_W, Inches(0.45),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    assumptions = [
        "Pricing mix: 20% Starter / 60% Professional / 20% Enterprise",
        "Y1–Y2 growth driven by private bank pilots and referrals",
        "Y3+ acceleration via PSB tenders and SI channel partnerships",
        "Y4–Y5 includes Southeast Asia revenue contribution (~15% of total)",
        "Headcount: 8 FTE (Y1) → 45 FTE (Y3) → 120 FTE (Y5)",
    ]
    add_multiline_text_box(slide, ["• " + a for a in assumptions],
                           MARGIN, Inches(4.8), CONTENT_W, Inches(2.1),
                           font_size=12, font_color=DARK_GRAY)

    # -----------------------------------------------------------------------
    # Slide 22 — Risk Analysis
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Risk Analysis & Mitigations")
    add_watermark(slide)
    add_slide_number(slide, 22)

    risk_data = [
        ["Risk",                          "Likelihood", "Impact", "Mitigation"],
        ["Long procurement cycles (PSBs)", "High",      "Medium", "SI partnerships (TCS, Infosys, Wipro) + start with private banks"],
        ["Legacy ATM hardware (XFS/SNMP)", "High",      "Medium", "XFS standard adapters + SNMP bridge layer in PayGuard agent"],
        ["Data localisation requirements", "Medium",    "High",   "On-premise deployment mode + private cloud option"],
        ["UPI growth reducing ATM usage",  "Medium",    "Medium", "Expand to POS terminals, CDMs, payment kiosks"],
        ["Competitor response (NCR/Diebold)","Low",     "High",   "Speed to market + India-first moat + 12-18 month head start"],
        ["Talent acquisition (AI/Django)",  "Medium",   "Medium", "Remote-first team + IIT/NIT campus recruitment"],
    ]
    add_table(slide, risk_data,
              MARGIN, Inches(1.55), CONTENT_W, Inches(4.5),
              col_widths=[0.28, 0.12, 0.1, 0.50],
              body_font_size=10)

    add_text_box(slide,
                 "All risks are known and have clear mitigation paths — no existential unknowns",
                 MARGIN, Inches(6.2), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 23 — Regulatory Tailwinds
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "Regulatory Tailwinds — Compliance as a Sales Driver")
    add_watermark(slide)
    add_slide_number(slide, 23)

    regs23 = [
        ("RBI/2011-12/314",
         ["ATM Failure Refund Mandate",
          "5 working day refund obligation",
          "₹100/day/customer for delays",
          "PayGuard automates tracking"]),
        ("RBI Master Direction\n(ATM Operations)",
         ["Uptime reporting requirements",
          "Incident documentation mandate",
          "PayGuard auto-generates reports",
          "Reduces compliance overhead"]),
        ("DPDP Act 2023",
         ["Digital Personal Data Protection",
          "SHA-256 hashing mandatory",
          "Data minimisation required",
          "Built into PayGuard core"]),
        ("PCI DSS v4.0",
         ["Payment Card Industry standard",
          "Transaction security controls",
          "Fraud detection requirements",
          "PayGuard aligns out-of-the-box"]),
    ]

    box_w = Inches(2.85)
    box_h = Inches(2.7)
    for i, (title, lines) in enumerate(regs23):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (box_w + gap), Inches(1.55),
                        box_w, box_h)

    add_text_box(slide,
                 "\"Regulatory pressure is our sales team\"",
                 MARGIN, Inches(4.4), CONTENT_W, Inches(0.45),
                 font_size=18, bold=True, italic=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    add_text_box(slide,
                 "Banks MUST spend on compliance. PayGuard turns mandatory spend into competitive advantage.",
                 MARGIN, Inches(4.95), CONTENT_W, Inches(0.45),
                 font_size=14, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=DARK_GRAY)

    add_text_box(slide,
                 "Indian banking IT compliance budgets are ring-fenced from discretionary cuts — recession-resistant revenue",
                 MARGIN, Inches(5.5), CONTENT_W, Inches(0.45),
                 font_size=13, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Slide 24 — Investment Thesis
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_title(slide, "The Investment Thesis")
    add_watermark(slide)
    add_slide_number(slide, 24)

    pillars = [
        ("1. Large Underserved Market",
         "$220M addressable in India alone. $18.2B globally by 2028. No AI-first player today."),
        ("2. First-Mover Advantage",
         "Only platform combining AI monitoring + predictive failure + customer transparency portal. 12-18 month head start on building this moat."),
        ("3. RBI Regulatory Tailwind",
         "Compliance spend is non-discretionary. RBI mandates drive urgency that our sales team cannot manufacture."),
        ("4. Modern Architecture",
         "React 19 + Django 6 + FastAPI vs SNMP polling and legacy C++ stacks. 10x faster to iterate. 10x cheaper to scale."),
        ("5. India-First Moat",
         "Bilingual EN/HI, OTP auth, ₹ formatting, RBI-native compliance. Global players cannot replicate quickly without 2-3 year localisation effort."),
    ]

    pillar_h = Inches(0.9)
    pillar_gap = Inches(0.12)
    top = Inches(1.55)
    for i, (title, body) in enumerate(pillars):
        py = top + i * (pillar_h + pillar_gap)
        p_box = slide.shapes.add_shape(1, MARGIN, py, CONTENT_W, pillar_h)
        p_box.fill.solid()
        p_box.fill.fore_color.rgb = PALE_GRAY if i % 2 == 0 else WHITE
        p_box.line.color.rgb = BLACK
        p_box.line.width = Pt(0.75)

        # Title portion (left 3.5 inches)
        add_text_box(slide, title,
                     MARGIN + Inches(0.15), py + Inches(0.1),
                     Inches(3.2), pillar_h - Inches(0.15),
                     font_size=13, bold=True, font_color=BLACK)

        # Body portion
        add_text_box(slide, body,
                     MARGIN + Inches(3.5), py + Inches(0.1),
                     CONTENT_W - Inches(3.65), pillar_h - Inches(0.15),
                     font_size=12, bold=False, font_color=DARK_GRAY)

    add_text_box(slide,
                 "PayGuard is building the Bloomberg Terminal of ATM infrastructure — mission-critical, sticky, and defensible.",
                 MARGIN, Inches(6.35), CONTENT_W, Inches(0.4),
                 font_size=13, bold=True, italic=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # -----------------------------------------------------------------------
    # Slide 25 — Contact / Next Steps
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_watermark(slide)
    add_slide_number(slide, 25)

    add_text_box(slide, "Ready to Pilot",
                 MARGIN, Inches(0.5), CONTENT_W, Inches(1.0),
                 font_size=48, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    # Horizontal rule
    rule25 = slide.shapes.add_shape(1, Inches(2), Inches(1.55), Inches(9.33), Pt(2))
    rule25.fill.solid()
    rule25.fill.fore_color.rgb = BLACK
    rule25.line.fill.background()

    add_text_box(slide, "We are seeking a pilot partner from India's private banking sector.",
                 MARGIN, Inches(1.75), CONTENT_W, Inches(0.45),
                 font_size=16, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=DARK_GRAY)

    next_steps = [
        ("Technical Demo",
         ["Live PayGuard dashboard walkthrough", "AI classifier + prediction demo", "PayGuard portal simulation", "Duration: 90 minutes"]),
        ("Pilot Agreement",
         ["500–1,000 ATM pilot scope", "3-month evaluation period", "Success metrics defined upfront", "Zero commitment post-pilot"]),
        ("Bank Introduction",
         ["Warm introduction to", "potential pilot bank partners", "Reference customers available", "NDA on request"]),
    ]

    col_w3 = Inches(3.9)
    col_h3 = Inches(2.6)
    col_top3 = Inches(2.35)
    col_gap3 = Inches(0.28)
    for i, (title, lines) in enumerate(next_steps):
        add_feature_box(slide, title, lines,
                        MARGIN + i * (col_w3 + col_gap3), col_top3,
                        col_w3, col_h3)

    # Company info
    add_text_box(slide,
                 "PayGuard  |  Product Intelligence Division",
                 MARGIN, Inches(5.2), CONTENT_W, Inches(0.4),
                 font_size=16, bold=True,
                 alignment=PP_ALIGN.CENTER, font_color=BLACK)

    add_text_box(slide,
                 "March 2026",
                 MARGIN, Inches(5.65), CONTENT_W, Inches(0.35),
                 font_size=14, bold=False,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    rule25b = slide.shapes.add_shape(1, Inches(3), Inches(6.1), Inches(7.33), Pt(1))
    rule25b.fill.solid()
    rule25b.fill.fore_color.rgb = LIGHT_GRAY
    rule25b.line.fill.background()

    add_text_box(slide,
                 "PayGuard — Transforming India's ATM Infrastructure through AI, Transparency, and Automation",
                 MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
                 font_size=13, bold=False, italic=True,
                 alignment=PP_ALIGN.CENTER, font_color=MID_GRAY)

    # -----------------------------------------------------------------------
    # Save
    # -----------------------------------------------------------------------
    output_path = "/Users/praanesh/PULSE/payguard_pitch_deck.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    build_deck()
