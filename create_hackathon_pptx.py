"""
PayGuard Hackathon Pitch Deck Generator
Produces: payguard_hackathon_deck.pptx
Color scheme: Strict black and white
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Color constants ──────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
MID_GRAY    = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY   = RGBColor(0x80, 0x80, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.47)   # ≈1.2 cm


# ── Presentation setup ───────────────────────────────────────────────────────
def make_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ─────────────────────────────────────────────────────────
def set_white_background(slide):
    """Fill slide background white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _set_cell_fill(cell, rgb):
    """Set a table cell's fill colour."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    # RGBColor stores as int tuple — access via index or unpack
    r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    srgbClr.set('val', '%02X%02X%02X' % (r, g, b))


def _set_shape_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_shape_border(shape, rgb=BLACK, width_pt=1):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


# ── Core builder functions ────────────────────────────────────────────────────
def add_text_box(slide, text, left, top, width, height,
                 font_size=16, bold=False,
                 align=PP_ALIGN.LEFT,
                 color=BLACK,
                 wrap=True,
                 italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = 'Calibri'
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_slide_title(slide, text):
    """Bold 30 pt title at top + thin black rule below."""
    title_box = add_text_box(
        slide, text,
        left=MARGIN, top=Inches(0.25),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.55),
        font_size=30, bold=True, align=PP_ALIGN.LEFT
    )
    # Thin black rule
    rule = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        MARGIN, Inches(0.85),
        SLIDE_W - 2 * MARGIN, Pt(1.5)
    )
    _set_shape_fill(rule, BLACK)
    rule.line.fill.background()
    return title_box


def add_slide_number(slide, num, total=12):
    add_text_box(
        slide,
        f"{num} / {total}",
        left=SLIDE_W - Inches(1.1), top=SLIDE_H - Inches(0.35),
        width=Inches(1.0), height=Inches(0.3),
        font_size=10, align=PP_ALIGN.RIGHT, color=DARK_GRAY
    )


def add_bottom_bar(slide, text):
    """Thin light-gray bar at bottom with centered small text."""
    bar = slide.shapes.add_shape(
        1,
        MARGIN, SLIDE_H - Inches(0.5),
        SLIDE_W - 2 * MARGIN, Inches(0.35)
    )
    _set_shape_fill(bar, LIGHT_GRAY)
    bar.line.fill.background()

    add_text_box(
        slide, text,
        left=MARGIN, top=SLIDE_H - Inches(0.5),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.35),
        font_size=10, align=PP_ALIGN.CENTER, color=BLACK
    )


def add_stat_box(slide, stat, label, left, top, width, height):
    """Gray fill box — large stat number + smaller label."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    _set_shape_fill(box, LIGHT_GRAY)
    _set_shape_border(box, BLACK, 1)

    # stat number
    add_text_box(
        slide, stat,
        left=left + Inches(0.1), top=top + Inches(0.15),
        width=width - Inches(0.2), height=height * 0.52,
        font_size=34, bold=True, align=PP_ALIGN.CENTER
    )
    # label
    add_text_box(
        slide, label,
        left=left + Inches(0.1), top=top + height * 0.52,
        width=width - Inches(0.2), height=height * 0.42,
        font_size=13, bold=False, align=PP_ALIGN.CENTER, color=DARK_GRAY
    )


def add_feature_box(slide, title, body, left, top, width, height):
    """Gray fill box — bold title + body text."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    _set_shape_fill(box, LIGHT_GRAY)
    _set_shape_border(box, BLACK, 1)

    title_h = Inches(0.38)
    add_text_box(
        slide, title,
        left=left + Inches(0.12), top=top + Inches(0.1),
        width=width - Inches(0.24), height=title_h,
        font_size=14, bold=True, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide, body,
        left=left + Inches(0.12), top=top + Inches(0.1) + title_h,
        width=width - Inches(0.24), height=height - title_h - Inches(0.2),
        font_size=12, bold=False, align=PP_ALIGN.LEFT, color=BLACK
    )


def add_styled_table(slide, headers, rows, left, top, width, height):
    """Black header row, alternating white / light-gray rows."""
    col_count = len(headers)
    row_count  = len(rows) + 1   # +1 for header

    tbl = slide.shapes.add_table(
        row_count, col_count,
        left, top, width, height
    ).table

    # Column widths — equal distribution
    col_w = width // col_count
    for i, col in enumerate(tbl.columns):
        col.width = col_w

    # Row height
    row_h = height // row_count
    for r in tbl.rows:
        r.height = row_h

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        _set_cell_fill(cell, BLACK)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = Pt(13)
        run.font.color.rgb = WHITE
        run.font.name  = 'Calibri'

    # Data rows
    for ri, row_data in enumerate(rows):
        fill_color = WHITE if ri % 2 == 0 else MID_GRAY
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            _set_cell_fill(cell, fill_color)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.size  = Pt(12)
            run.font.color.rgb = BLACK
            run.font.name  = 'Calibri'


def add_arch_box(slide, text, left, top, width, height, font_size=13):
    """Box styled for architecture diagram."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    _set_shape_fill(box, LIGHT_GRAY)
    _set_shape_border(box, BLACK, 1.2)
    add_text_box(
        slide, text,
        left=left + Inches(0.08), top=top + Inches(0.05),
        width=width - Inches(0.16), height=height - Inches(0.1),
        font_size=font_size, bold=False, align=PP_ALIGN.LEFT
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)

    # Big product name
    add_text_box(
        slide, "PayGuard",
        left=MARGIN, top=Inches(1.4),
        width=SLIDE_W - 2 * MARGIN, height=Inches(1.4),
        font_size=72, bold=True, align=PP_ALIGN.CENTER
    )

    # Subtitle
    add_text_box(
        slide,
        "Real-Time ATM Monitoring + Customer Transparency Platform",
        left=MARGIN, top=Inches(2.85),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.6),
        font_size=22, bold=False, align=PP_ALIGN.CENTER, color=DARK_GRAY
    )

    # Tagline
    add_text_box(
        slide,
        '"When your ATM fails, you shouldn\'t be left in the dark."',
        left=MARGIN, top=Inches(3.55),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.5),
        font_size=18, bold=False, align=PP_ALIGN.CENTER, italic=True
    )

    # Hackathon + date
    add_text_box(
        slide, "Hackathon Submission  |  March 2026",
        left=MARGIN, top=Inches(4.2),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.4),
        font_size=14, bold=False, align=PP_ALIGN.CENTER, color=DARK_GRAY
    )

    # Tech stack pills — simulate with one text box, spaced items
    pills_text = "Django 6    |    React 19    |    FastAPI    |    Django Channels    |    Redis"
    pill_box = slide.shapes.add_shape(
        1,
        MARGIN, Inches(4.85),
        SLIDE_W - 2 * MARGIN, Inches(0.48)
    )
    _set_shape_fill(pill_box, LIGHT_GRAY)
    _set_shape_border(pill_box, BLACK, 1)

    add_text_box(
        slide, pills_text,
        left=MARGIN, top=Inches(4.87),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.44),
        font_size=13, bold=True, align=PP_ALIGN.CENTER
    )

    # Thin rule above tech stack
    rule = slide.shapes.add_shape(1, MARGIN, Inches(4.78), SLIDE_W - 2 * MARGIN, Pt(1))
    _set_shape_fill(rule, BLACK)
    rule.line.fill.background()


def slide_02_problem_scale(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "52.5 Million Failed ATM Transactions. Every Month.")

    # 2x2 stat boxes
    box_w = Inches(2.9)
    box_h = Inches(2.0)
    gap   = Inches(0.28)
    start_x = MARGIN
    start_y = Inches(1.1)

    stats = [
        ("255,000",      "ATMs in India's network"),
        ("1.05B/month",  "Transactions processed"),
        ("5–8%",         "Failure rate = 52.5–84M failures/month"),
        ("70%",          "Failures undetected for >30 minutes"),
    ]

    positions = [
        (start_x,               start_y),
        (start_x + box_w + gap, start_y),
        (start_x,               start_y + box_h + gap),
        (start_x + box_w + gap, start_y + box_h + gap),
    ]

    for (stat, label), (lft, tp) in zip(stats, positions):
        add_stat_box(slide, stat, label, lft, tp, box_w, box_h)

    # Right side — additional context block
    ctx_x = start_x + 2 * (box_w + gap) + Inches(0.1)
    ctx_w = SLIDE_W - ctx_x - MARGIN
    ctx_box = slide.shapes.add_shape(1, ctx_x, start_y, ctx_w, 2 * box_h + gap)
    _set_shape_fill(ctx_box, LIGHT_GRAY)
    _set_shape_border(ctx_box, BLACK, 1)

    context_text = (
        "ATM Uptime — India vs World\n\n"
        "India today:   90–93%\n"
        "Developed markets:  98%+\n\n"
        "That 5–8% gap costs customers\n"
        "millions of failed transactions\n"
        "and hours of uncertainty daily.\n\n"
        "Rs. 630 crore in ATM fraud\n"
        "reported in FY 2022-23 (RBI)"
    )
    add_text_box(
        slide, context_text,
        left=ctx_x + Inches(0.15), top=start_y + Inches(0.15),
        width=ctx_w - Inches(0.3), height=2 * box_h + gap - Inches(0.3),
        font_size=13, align=PP_ALIGN.LEFT
    )

    add_bottom_bar(slide, "ATM uptime in India: 90–93% vs 98%+ in developed markets  |  Source: RBI, NPCI 2023")
    add_slide_number(slide, 2)


def slide_03_problem_cx(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "What Happens When Your ATM Transaction Fails?")

    # Central statement
    add_text_box(
        slide,
        "Your money is gone. Your card is stuck. The screen is blank.",
        left=MARGIN, top=Inches(1.05),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.65),
        font_size=22, bold=True, align=PP_ALIGN.CENTER
    )

    # 4 pain boxes
    pains = [
        ("No real-time status",     "Is my money actually debited?"),
        ("No refund tracking",      "When will I get my money back?"),
        ("No engineer visibility",  "Is anyone even fixing this?"),
        ("Helpline: 45-min wait",   "The only option available today"),
    ]

    box_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 4
    box_h = Inches(2.5)
    gap   = Inches(0.2)
    top_y = Inches(1.85)

    for i, (title, body) in enumerate(pains):
        lft = MARGIN + i * (box_w + gap)
        add_feature_box(slide, title, body, lft, top_y, box_w, box_h)

    add_bottom_bar(
        slide,
        "73% of Indian ATM users have experienced a failed transaction  (NPCI 2023)"
    )
    add_slide_number(slide, 3)


def slide_04_what_we_built(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "PayGuard: End-to-End ATM Intelligence")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    col_h = Inches(4.8)
    top_y = Inches(1.05)

    # Left column — operator
    left_col = slide.shapes.add_shape(1, MARGIN, top_y, col_w, col_h)
    _set_shape_fill(left_col, LIGHT_GRAY)
    _set_shape_border(left_col, BLACK, 1)

    add_text_box(
        slide, "For Bank Operators — PULSE Engine",
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.45),
        font_size=15, bold=True, align=PP_ALIGN.LEFT
    )

    op_features = (
        "  Real-time fleet monitoring\n"
        "  AI root cause classification\n"
        "  Predictive failure detection\n"
        "  Autonomous self-healing\n"
        "  Fraud detection engine\n"
        "  Incident management & audit log"
    )
    add_text_box(
        slide, op_features,
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.65),
        width=col_w - Inches(0.24), height=col_h - Inches(0.8),
        font_size=14, align=PP_ALIGN.LEFT
    )

    # Right column — customer portal
    right_x = MARGIN + col_w + Inches(0.4)
    right_col = slide.shapes.add_shape(1, right_x, top_y, col_w, col_h)
    _set_shape_fill(right_col, LIGHT_GRAY)
    _set_shape_border(right_col, BLACK, 1)

    add_text_box(
        slide, "For Customers — PayGuard Portal",
        left=right_x + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.45),
        font_size=15, bold=True, align=PP_ALIGN.LEFT
    )

    cx_features = (
        "  Real-time failed transaction status\n"
        "  Refund tracking with ETA\n"
        "  OTP login (no password / no app)\n"
        "  Bilingual English + Hindi\n"
        "  Live WebSocket updates\n"
        "  RBI compliance messaging baked in"
    )
    add_text_box(
        slide, cx_features,
        left=right_x + Inches(0.12), top=top_y + Inches(0.65),
        width=col_w - Inches(0.24), height=col_h - Inches(0.8),
        font_size=14, align=PP_ALIGN.LEFT
    )

    add_bottom_bar(slide, "Built in one hackathon. Production-ready architecture.")
    add_slide_number(slide, 4)


def slide_05_architecture(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "Three-Service Architecture")

    top_y = Inches(1.05)

    # ── Tier 1: Frontend ──
    fe_w = Inches(5.5)
    fe_h = Inches(0.9)
    fe_x = (SLIDE_W - fe_w) / 2
    fe_y = top_y

    fe_box = slide.shapes.add_shape(1, fe_x, fe_y, fe_w, fe_h)
    _set_shape_fill(fe_box, LIGHT_GRAY)
    _set_shape_border(fe_box, BLACK, 1.2)

    add_text_box(
        slide,
        "React 19 Frontend\nDashboard + Customer Portal  (TypeScript + Vite + Redux)",
        left=fe_x + Inches(0.12), top=fe_y + Inches(0.05),
        width=fe_w - Inches(0.24), height=fe_h - Inches(0.1),
        font_size=13, bold=False, align=PP_ALIGN.CENTER
    )

    # Arrow down from frontend
    arr1_x = fe_x + fe_w / 2 - Pt(4)
    arr1_y = fe_y + fe_h
    arr1_h = Inches(0.38)
    arr1 = slide.shapes.add_shape(1, arr1_x, arr1_y, Pt(8), arr1_h)
    _set_shape_fill(arr1, BLACK)
    arr1.line.fill.background()

    add_text_box(
        slide, "REST / WebSocket",
        left=arr1_x + Inches(0.15), top=arr1_y + Inches(0.05),
        width=Inches(1.5), height=Inches(0.3),
        font_size=10, color=DARK_GRAY
    )

    # ── Tier 2: Django backend ──
    dj_w = Inches(7.5)
    dj_h = Inches(1.4)
    dj_x = (SLIDE_W - dj_w) / 2
    dj_y = fe_y + fe_h + arr1_h

    dj_box = slide.shapes.add_shape(1, dj_x, dj_y, dj_w, dj_h)
    _set_shape_fill(dj_box, LIGHT_GRAY)
    _set_shape_border(dj_box, BLACK, 1.2)

    dj_text = (
        "Django 6 REST API + Django Channels  (port 8000)\n"
        "40+ REST endpoints  |  Django Channels WebSocket  |  SQLite + 16 data models\n"
        "JWT authentication  |  Redis channel layer  |  Celery-ready async views"
    )
    add_text_box(
        slide, dj_text,
        left=dj_x + Inches(0.15), top=dj_y + Inches(0.1),
        width=dj_w - Inches(0.3), height=dj_h - Inches(0.15),
        font_size=13, align=PP_ALIGN.CENTER
    )

    # Arrow down from Django
    arr2_x = dj_x + dj_w / 2 - Pt(4)
    arr2_y = dj_y + dj_h
    arr2_h = Inches(0.38)
    arr2 = slide.shapes.add_shape(1, arr2_x, arr2_y, Pt(8), arr2_h)
    _set_shape_fill(arr2, BLACK)
    arr2.line.fill.background()

    add_text_box(
        slide, "HTTP proxy",
        left=arr2_x + Inches(0.15), top=arr2_y + Inches(0.05),
        width=Inches(1.2), height=Inches(0.3),
        font_size=10, color=DARK_GRAY
    )

    # ── Tier 3: FastAPI AI service ──
    ai_w = Inches(7.5)
    ai_h = Inches(1.4)
    ai_x = (SLIDE_W - ai_w) / 2
    ai_y = dj_y + dj_h + arr2_h

    ai_box = slide.shapes.add_shape(1, ai_x, ai_y, ai_w, ai_h)
    _set_shape_fill(ai_box, LIGHT_GRAY)
    _set_shape_border(ai_box, BLACK, 1.2)

    ai_text = (
        "FastAPI AI Microservice  (port 8001)\n"
        "Failure Classifier  |  Predictive Model  |  Anomaly Detector  |  Fraud Engine\n"
        "NumPy statistical models  |  Sub-second inference  |  Async endpoints"
    )
    add_text_box(
        slide, ai_text,
        left=ai_x + Inches(0.15), top=ai_y + Inches(0.1),
        width=ai_w - Inches(0.3), height=ai_h - Inches(0.15),
        font_size=13, align=PP_ALIGN.CENTER
    )

    # ── Tech badge row ──
    badges = ["React 19 + TypeScript", "Django 6 + DRF", "FastAPI + NumPy", "Django Channels"]
    badge_w = Inches(2.7)
    badge_h = Inches(0.38)
    badge_gap = Inches(0.22)
    total_badge_w = len(badges) * badge_w + (len(badges) - 1) * badge_gap
    badge_x0 = (SLIDE_W - total_badge_w) / 2
    badge_y  = ai_y + ai_h + Inches(0.25)

    for i, badge_text in enumerate(badges):
        bx = badge_x0 + i * (badge_w + badge_gap)
        b = slide.shapes.add_shape(1, bx, badge_y, badge_w, badge_h)
        _set_shape_fill(b, BLACK)
        b.line.fill.background()
        add_text_box(
            slide, badge_text,
            left=bx, top=badge_y,
            width=badge_w, height=badge_h,
            font_size=12, bold=True, align=PP_ALIGN.CENTER, color=WHITE
        )

    add_slide_number(slide, 5)


def slide_06_ai_engine(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "Four AI Models, One Pipeline")

    boxes_data = [
        (
            "1. Failure Classifier",
            "eventCode → category + confidence + self-heal action\n\n"
            "Lookup table with confidence scoring.\n"
            "7 known event codes mapped to 5 failure categories:\n"
            "NETWORK / CASH_JAM / HARDWARE / SERVER / FRAUD\n\n"
            "Output: category, confidence, recommended action."
        ),
        (
            "2. Predictive Model",
            "Health time-series slope analysis\n→ P(failure) + hours to failure\n\n"
            "Runs rolling mean, variance, and slope on the last N health snapshots.\n\n"
            "Catches declining ATMs before they fail completely."
        ),
        (
            "3. Anomaly Detector",
            "Z-score on error rate  (threshold: σ > 3.0)\n\n"
            "Detects:\n"
            "  RAPID_FAILURES\n"
            "  MALWARE_PATTERN\n"
            "  CARD_SKIMMING\n\n"
            "Compares live error rate against rolling historical baseline."
        ),
        (
            "4. Fraud Detector",
            "3-heuristic ensemble — sub-second per transaction:\n\n"
            "  Velocity: 3+ withdrawals in 10 min (same card)\n"
            "  Amount: Z-score > 3.5σ from card baseline\n"
            "  Travel: Haversine >100 km in <60 min\n\n"
            "Confidence scales with heuristic count."
        ),
    ]

    box_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 4
    box_h = Inches(4.7)
    gap   = Inches(0.2)
    top_y = Inches(1.05)

    for i, (title, body) in enumerate(boxes_data):
        lft = MARGIN + i * (box_w + gap)
        add_feature_box(slide, title, body, lft, top_y, box_w, box_h)

    add_bottom_bar(
        slide,
        "All 4 models run as async services.  Self-healing triggers automatically for 4 of 6 failure categories."
    )
    add_slide_number(slide, 6)


def slide_07_self_healing(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "Autonomous Recovery in Under 2 Minutes")

    headers = ["Root Cause", "Self-Heal Action", "Auto?", "Typical MTTR"]
    rows = [
        ["NETWORK",         "Switch to backup network path",      "Yes",  "< 90 seconds"],
        ["SERVER / SWITCH", "Restart service",                    "Yes",  "< 2 minutes"],
        ["TIMEOUT",         "Flush cache + retry",                "Yes",  "< 60 seconds"],
        ["CASH_JAM",        "Dispatch engineer alert",            "No",   "30–120 minutes"],
        ["FRAUD",           "Freeze ATM instantly",               "No",   "Immediate"],
        ["HARDWARE",        "Alert engineer + open incident",     "No",   "30–120 minutes"],
    ]

    add_styled_table(
        slide, headers, rows,
        left=MARGIN, top=Inches(1.05),
        width=SLIDE_W - 2 * MARGIN, height=Inches(5.0)
    )

    add_bottom_bar(
        slide,
        "Industry average MTTR: 4+ hours.  PayGuard auto-resolves 4 of 6 failure categories."
    )
    add_slide_number(slide, 7)


def slide_08_customer_portal(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "The First Real-Time ATM Incident Portal in India")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    top_y = Inches(1.05)
    col_h = Inches(4.5)

    # Left — flow
    left_box = slide.shapes.add_shape(1, MARGIN, top_y, col_w, col_h)
    _set_shape_fill(left_box, LIGHT_GRAY)
    _set_shape_border(left_box, BLACK, 1)

    add_text_box(
        slide, "Customer Flow",
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.42),
        font_size=15, bold=True
    )

    flow_text = (
        "1.  Customer enters phone number\n\n"
        "2.  OTP sent (no account required)\n\n"
        "3.  View all failed transactions\n\n"
        "4.  Live status updates via WebSocket\n     (no manual refresh needed)\n\n"
        "5.  Refund ETA displayed prominently\n\n"
        "6.  SMS deep-link — works without app"
    )
    add_text_box(
        slide, flow_text,
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.62),
        width=col_w - Inches(0.24), height=col_h - Inches(0.75),
        font_size=13
    )

    # Right — design decisions
    right_x = MARGIN + col_w + Inches(0.4)
    right_box = slide.shapes.add_shape(1, right_x, top_y, col_w, col_h)
    _set_shape_fill(right_box, LIGHT_GRAY)
    _set_shape_border(right_box, BLACK, 1)

    add_text_box(
        slide, "Key Design Decisions",
        left=right_x + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.42),
        font_size=15, bold=True
    )

    decisions_text = (
        "OTP auth — zero setup, familiar from UPI\n\n"
        "SHA-256 phone hashing — DPDP Act 2023\n"
        "compliant, no raw PII stored\n\n"
        "WebSocket + exponential backoff — built\n"
        "for Indian 2G / 3G mobile networks\n\n"
        "SMS deep-link — no smartphone app needed\n\n"
        '"Your money is safe" — first thing\n'
        "the customer sees on every screen"
    )
    add_text_box(
        slide, decisions_text,
        left=right_x + Inches(0.12), top=top_y + Inches(0.62),
        width=col_w - Inches(0.24), height=col_h - Inches(0.75),
        font_size=13
    )

    # Status flow bar at bottom
    statuses = ["DETECTED", "INVESTIGATING", "ENGINEER\nDISPATCHED", "RESOLVING", "REFUND\nINITIATED", "RESOLVED"]
    status_w = (SLIDE_W - 2 * MARGIN - Inches(0.5)) / len(statuses)
    status_h = Inches(0.55)
    status_y = top_y + col_h + Inches(0.18)
    gap_s = Inches(0.1)

    for i, s in enumerate(statuses):
        sx = MARGIN + i * (status_w + gap_s)
        sb = slide.shapes.add_shape(1, sx, status_y, status_w, status_h)
        _set_shape_fill(sb, BLACK if i == 0 else LIGHT_GRAY)
        _set_shape_border(sb, BLACK, 1)
        add_text_box(
            slide, s,
            left=sx, top=status_y,
            width=status_w, height=status_h,
            font_size=9, bold=(i == 0), align=PP_ALIGN.CENTER,
            color=WHITE if i == 0 else BLACK
        )

    add_slide_number(slide, 8)


def slide_09_bilingual(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "English + Hindi. Not a Translation. A Native Experience.")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    top_y = Inches(1.1)
    col_h = Inches(2.9)

    # English box
    en_box = slide.shapes.add_shape(1, MARGIN, top_y, col_w, col_h)
    _set_shape_fill(en_box, LIGHT_GRAY)
    _set_shape_border(en_box, BLACK, 1)

    add_text_box(
        slide, "English",
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.38),
        font_size=14, bold=True
    )

    en_text = (
        "Cash jam — your \u20b93,000 was debited but cash was not dispensed.\n\n"
        "Your money is safe and will be refunded within 5 working days "
        "as per RBI guidelines.\n\n"
        "Reference: TXN-2024-00418\n"
        "Engineer dispatched. ETA: 45 minutes."
    )
    add_text_box(
        slide, en_text,
        left=MARGIN + Inches(0.12), top=top_y + Inches(0.55),
        width=col_w - Inches(0.24), height=col_h - Inches(0.65),
        font_size=13
    )

    # Hindi box
    hi_x = MARGIN + col_w + Inches(0.4)
    hi_box = slide.shapes.add_shape(1, hi_x, top_y, col_w, col_h)
    _set_shape_fill(hi_box, LIGHT_GRAY)
    _set_shape_border(hi_box, BLACK, 1)

    add_text_box(
        slide, "Hindi (\u0939\u093f\u0928\u094d\u0926\u0940)",
        left=hi_x + Inches(0.12), top=top_y + Inches(0.1),
        width=col_w - Inches(0.24), height=Inches(0.38),
        font_size=14, bold=True
    )

    hi_text = (
        "\u0915\u0948\u0936 \u091c\u093e\u092e \u2014 \u0906\u092a\u0915\u093e "
        "\u20b93,000 \u0921\u0947\u092c\u093f\u091f \u0939\u094b \u0917\u092f\u093e "
        "\u0932\u0947\u0915\u093f\u0928 \u0928\u0915\u0926 \u0928\u0939\u0940\u0902 "
        "\u0928\u093f\u0915\u0932\u093e\u0964\n\n"
        "RBI \u0928\u093f\u092f\u092e\u094b\u0902 \u0915\u0947 \u0905\u0928\u0941\u0938\u093e\u0930 "
        "\u0906\u092a\u0915\u093e \u092a\u0948\u0938\u093e 5 \u0915\u093e\u0930\u094d\u092f "
        "\u0926\u093f\u0935\u0938\u094b\u0902 \u092e\u0947\u0902 \u0935\u093e\u092a\u0938 "
        "\u0915\u093f\u092f\u093e \u091c\u093e\u090f\u0917\u093e\u0964\n\n"
        "\u0938\u0902\u0926\u0930\u094d\u092d: TXN-2024-00418\n"
        "\u0907\u0902\u091c\u0940\u0928\u093f\u092f\u0930 \u092d\u0947\u091c\u093e \u0917\u092f\u093e\u0964 "
        "ETA: 45 \u092e\u093f\u0928\u091f"
    )
    add_text_box(
        slide, hi_text,
        left=hi_x + Inches(0.12), top=top_y + Inches(0.55),
        width=col_w - Inches(0.24), height=col_h - Inches(0.65),
        font_size=13
    )

    # Coverage note
    add_text_box(
        slide,
        "Every status label, error message, failure reason, engineer update, and UI string has both EN and HI variants.",
        left=MARGIN, top=top_y + col_h + Inches(0.22),
        width=SLIDE_W - 2 * MARGIN, height=Inches(0.4),
        font_size=13, align=PP_ALIGN.CENTER, italic=True
    )

    add_bottom_bar(
        slide,
        "65% of India's population primarily uses Hindi. We built for them."
    )
    add_slide_number(slide, 9)


def slide_10_fraud(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "3-Heuristic Fraud Ensemble")

    boxes_data = [
        (
            "1. Velocity Check",
            "3+ withdrawals from the same card\nin a 10-minute rolling window\n\n"
            "Confidence scales with transaction count.\n\n"
            "Runs per transaction against card history.\n\n"
            "Flags: RAPID_FAILURES\n\n"
            "Trigger: AnomalyFlag created,\nengineer alerted."
        ),
        (
            "2. Amount Anomaly",
            "Z-score vs card's historical\ntransaction baseline\n\n"
            "Threshold: |Z| > 3.5\u03c3\n\n"
            "Each card builds its own statistical\nprofile over time.\n\n"
            "Flags: UNUSUAL_WITHDRAWAL\n\n"
            "Accounts for card-holder-specific\nspending patterns."
        ),
        (
            "3. Impossible Travel",
            "Haversine distance between\nconsecutive transactions\n\n"
            "Threshold: >100 km in <60 min\n\n"
            "Catches cloned cards used\nsimultaneously at distance.\n\n"
            "Flags: geographic impossibility\n\n"
            "Uses exact lat/long of each ATM."
        ),
    ]

    box_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 3
    box_h = Inches(4.7)
    gap   = Inches(0.2)
    top_y = Inches(1.05)

    for i, (title, body) in enumerate(boxes_data):
        lft = MARGIN + i * (box_w + gap)
        add_feature_box(slide, title, body, lft, top_y, box_w, box_h)

    add_bottom_bar(
        slide,
        "All 3 heuristics run as background threads.  Result: AnomalyFlag created, ATM status updated, engineer alerted — in under 1 second."
    )
    add_slide_number(slide, 10)


def slide_11_demo(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "What You'll See in the Demo")

    steps = [
        (
            "1. Trigger a Failure",
            "POST /api/logs/ingest/ with event code CASH_DISPENSE_FAIL.\n"
            "Raw log entry created in the database."
        ),
        (
            "2. Watch the Pipeline",
            "Log saved  \u2192  AI Classifier fires  \u2192  Incident created as CASH_JAM\n"
            "\u2192  Health score recalculated  \u2192  Self-heal action selected\n"
            "\u2192  Dashboard updates via WebSocket  (no page refresh)"
        ),
        (
            "3. Customer Portal",
            "Enter phone  \u2192  OTP  \u2192  See failed transaction\n"
            "Status: DETECTED  \u2192  live advance to INVESTIGATING\n"
            "\u2192  ENGINEER DISPATCHED  (all via WebSocket push)"
        ),
        (
            "4. Fraud Scenario",
            "3 rapid withdrawals from same card  (within 60 seconds)\n"
            "\u2192  Velocity heuristic fires  \u2192  RAPID_FAILURES flag\n"
            "\u2192  ATM status: UNDER_REVIEW  \u2192  AnomalyFlag on dashboard"
        ),
    ]

    box_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 2
    box_h = Inches(2.25)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    top_y = Inches(1.05)

    positions = [
        (MARGIN,              top_y),
        (MARGIN + box_w + gap_x, top_y),
        (MARGIN,              top_y + box_h + gap_y),
        (MARGIN + box_w + gap_x, top_y + box_h + gap_y),
    ]

    for (title, body), (lft, tp) in zip(steps, positions):
        add_feature_box(slide, title, body, lft, tp, box_w, box_h)

    add_bottom_bar(slide, "Fully working end-to-end.  No mocks.  No slides pretending to be a product.")
    add_slide_number(slide, 11)


def slide_12_differentiators(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_white_background(slide)
    add_slide_title(slide, "Six Things No One Else Has Built")

    items = [
        ("1", "Customer-facing real-time ATM incident portal",
               "First of its kind in India. Banks have internal dashboards. Customers have nothing. Until now."),
        ("2", "Impossible travel fraud detection (Haversine)",
               "Real-time Haversine distance check flags cloned cards used >100 km apart in <60 min."),
        ("3", "Autonomous self-healing for 4 failure categories",
               "SWITCH_NETWORK, RESTART_SERVICE, FLUSH_CACHE — auto-resolved in under 2 minutes."),
        ("4", "Bilingual WebSocket portal with exponential backoff",
               "English + Hindi live updates designed for Indian 2G/3G network conditions."),
        ("5", "RBI + DPDP Act 2023 compliance by design",
               "SHA-256 phone hashing, RBI refund timeline messaging, DPDP-compliant data model."),
        ("6", "Full stack in one hackathon",
               "AI + WebSocket + OTP auth + SMS deep-link + fraud detection + bilingual UX — all integrated."),
    ]

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    box_h = Inches(1.45)
    gap_x = Inches(0.3)
    gap_y = Inches(0.15)
    top_y = Inches(1.05)

    for i, (num, title, body) in enumerate(items):
        row = i // 2
        col = i % 2
        lft = MARGIN + col * (col_w + gap_x)
        tp  = top_y + row * (box_h + gap_y)

        bx = slide.shapes.add_shape(1, lft, tp, col_w, box_h)
        _set_shape_fill(bx, LIGHT_GRAY)
        _set_shape_border(bx, BLACK, 1)

        # Number badge
        num_w = Inches(0.38)
        nb = slide.shapes.add_shape(1, lft + Inches(0.12), tp + Inches(0.12), num_w, num_w)
        _set_shape_fill(nb, BLACK)
        nb.line.fill.background()
        add_text_box(
            slide, num,
            left=lft + Inches(0.12), top=tp + Inches(0.1),
            width=num_w, height=num_w,
            font_size=16, bold=True, align=PP_ALIGN.CENTER, color=WHITE
        )

        # Title
        add_text_box(
            slide, title,
            left=lft + Inches(0.62), top=tp + Inches(0.1),
            width=col_w - Inches(0.74), height=Inches(0.42),
            font_size=13, bold=True
        )
        # Body
        add_text_box(
            slide, body,
            left=lft + Inches(0.62), top=tp + Inches(0.52),
            width=col_w - Inches(0.74), height=box_h - Inches(0.62),
            font_size=11, color=BLACK
        )

    # Final footer bar — black
    footer_h = Inches(0.42)
    footer_y = SLIDE_H - footer_h - Inches(0.08)
    fb = slide.shapes.add_shape(1, MARGIN, footer_y, SLIDE_W - 2 * MARGIN, footer_h)
    _set_shape_fill(fb, BLACK)
    fb.line.fill.background()
    add_text_box(
        slide,
        "PayGuard  |  Django 6  \u00b7  React 19  \u00b7  FastAPI  \u00b7  Django Channels  |  March 2026",
        left=MARGIN, top=footer_y,
        width=SLIDE_W - 2 * MARGIN, height=footer_h,
        font_size=12, bold=True, align=PP_ALIGN.CENTER, color=WHITE
    )

    add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = make_presentation()

    slide_01_title(prs)
    slide_02_problem_scale(prs)
    slide_03_problem_cx(prs)
    slide_04_what_we_built(prs)
    slide_05_architecture(prs)
    slide_06_ai_engine(prs)
    slide_07_self_healing(prs)
    slide_08_customer_portal(prs)
    slide_09_bilingual(prs)
    slide_10_fraud(prs)
    slide_11_demo(prs)
    slide_12_differentiators(prs)

    output = "/Users/praanesh/PULSE/payguard_hackathon_deck.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
