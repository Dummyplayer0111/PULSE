"""
Generate PULSE hackathon presentation (15 slides).
Run: python make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BG      = RGBColor(0x0b, 0x0b, 0x0f)   # near-black background
ACCENT  = RGBColor(0x6c, 0x63, 0xff)   # purple accent
ACCENT2 = RGBColor(0x00, 0xd4, 0xff)   # cyan
GREEN   = RGBColor(0x10, 0xb9, 0x81)   # green
RED     = RGBColor(0xef, 0x44, 0x44)   # red
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)   # amber
WHITE   = RGBColor(0xff, 0xff, 0xff)
GREY    = RGBColor(0x94, 0xa3, 0xb8)
DARK_CARD = RGBColor(0x16, 0x16, 0x22) # card background

W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.5)     # 16:9 widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_bg(slide, color=BG):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_heading(slide, title, subtitle=None):
    """Standard slide header with purple accent bar."""
    # Accent bar
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)
    # Title
    add_textbox(slide, title,
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                    font_size=16, color=GREY)


def add_card(slide, x, y, w, h, title, body_lines, title_color=ACCENT2):
    """Dark card with title + bullet body."""
    add_rect(slide, x, y, w, h, DARK_CARD)
    # Card top accent
    add_rect(slide, x, y, w, Inches(0.04), title_color)
    # Title
    add_textbox(slide, title,
                x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.4),
                font_size=13, bold=True, color=title_color)
    # Body
    body_text = "\n".join(f"• {line}" for line in body_lines)
    add_textbox(slide, body_text,
                x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
                font_size=11, color=WHITE)


def add_badge(slide, text, x, y, color=ACCENT):
    """Small coloured pill/badge."""
    add_rect(slide, x, y, Inches(1.5), Inches(0.32), color)
    add_textbox(slide, text, x, y, Inches(1.5), Inches(0.32),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)

# Big gradient-like background blocks
add_rect(s, 0, 0, W * 0.55, H, RGBColor(0x10, 0x10, 0x1a))
add_rect(s, W * 0.55, 0, W * 0.45, H, RGBColor(0x0e, 0x0e, 0x18))

# Accent stripe
add_rect(s, 0, 0, Inches(0.12), H, ACCENT)

# PULSE title
add_textbox(s, "PULSE",
            Inches(0.5), Inches(1.5), Inches(7), Inches(1.6),
            font_size=80, bold=True, color=ACCENT)

# Full name
add_textbox(s, "Predictive Unified Log & Surveillance Engine",
            Inches(0.5), Inches(3.1), Inches(7.5), Inches(0.6),
            font_size=20, color=ACCENT2)

# Tagline
add_textbox(s, "AI-Powered ATM & Payment Channel Intelligence Platform",
            Inches(0.5), Inches(3.75), Inches(8), Inches(0.45),
            font_size=15, color=GREY)

# Right side visual description
add_textbox(s, "Real-time monitoring\nAI root cause analysis\nAutonomous self-healing\nMultilingual notifications",
            Inches(8.8), Inches(2.5), Inches(4), Inches(2.5),
            font_size=16, color=WHITE)

# Badges
for i, (label, col) in enumerate([("AI-Native", ACCENT), ("Real-Time", ACCENT2), ("Self-Healing", GREEN)]):
    add_badge(s, label, Inches(0.5 + i * 1.75), Inches(4.5), col)

# Footer
add_textbox(s, "Hackathon 2026",
            Inches(0.5), H - Inches(0.6), Inches(6), Inches(0.4),
            font_size=12, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "The Problem", "India's ATM network is failing silently — at massive scale")

# Big stat cards
stats = [
    ("2.5 Lakh+", "ATMs across India", RED),
    ("~30%", "Experience monthly faults", AMBER),
    ("Hours", "Average fault-to-fix time", RED),
    ("₹1000 Cr+", "Annual downtime losses", AMBER),
]
for i, (num, label, col) in enumerate(stats):
    cx = Inches(0.5 + i * 3.2)
    add_rect(s, cx, Inches(1.4), Inches(3.0), Inches(1.6), DARK_CARD)
    add_rect(s, cx, Inches(1.4), Inches(3.0), Inches(0.05), col)
    add_textbox(s, num, cx, Inches(1.5), Inches(3.0), Inches(0.8),
                font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, label, cx, Inches(2.3), Inches(3.0), Inches(0.5),
                font_size=12, color=GREY, align=PP_ALIGN.CENTER)

# Pain points
pain = [
    "No real-time visibility — operators learn of failures from customer complaints",
    "Manual log analysis — engineers sift through thousands of log lines daily",
    "Reactive, not predictive — faults are fixed AFTER customers are already impacted",
    "Language barrier — customer notifications sent in English only to non-English users",
    "Slow escalation — average 2–4 hours from fault detection to resolution",
    "No automated remediation — every fix requires a field engineer or manual restart",
]
add_textbox(s, "Current Pain Points", Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
            font_size=16, bold=True, color=ACCENT2)
for i, pt in enumerate(pain):
    row, col_idx = divmod(i, 3)
    add_textbox(s, f"✗  {pt}",
                Inches(0.5 + col_idx * 4.3), Inches(3.75 + row * 0.65),
                Inches(4.2), Inches(0.55),
                font_size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Market Context
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Market Context", "The scale of India's digital payment infrastructure")

items = [
    ("2.5 Lakh+ ATMs", "Deployed across urban & rural India (RBI 2024)"),
    ("₹37 Lakh Cr/year", "Cash withdrawn from ATMs annually"),
    ("UPI + Card transactions", "10 billion+ monthly transactions across payment channels"),
    ("100+ banks", "Operating independent ATM networks with no unified monitoring"),
    ("Multilingual users", "22 official languages — most bank apps support only English/Hindi"),
    ("Rural penetration", "40% of ATMs in tier-3 and below towns — highest failure risk"),
]
for i, (title, desc) in enumerate(items):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.4 + row * 1.65)
    add_rect(s, x, y, Inches(6.0), Inches(1.45), DARK_CARD)
    add_rect(s, x, y, Inches(0.06), Inches(1.45), ACCENT)
    add_textbox(s, title, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.45),
                font_size=16, bold=True, color=ACCENT2)
    add_textbox(s, desc, x + Inches(0.2), y + Inches(0.6), Inches(5.7), Inches(0.7),
                font_size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Solution Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Our Solution — PULSE", "Detect → Classify → Heal → Notify, all in < 2 seconds")

# Flow diagram
steps = [
    ("📥 Ingest", "ATM logs stream in via REST API in real time", ACCENT),
    ("🤖 Classify", "FastAPI AI classifies root cause with confidence score", ACCENT2),
    ("🚨 Incident", "Incident auto-created if confidence ≥ 65% on ERROR/CRITICAL", AMBER),
    ("🔧 Self-Heal", "Automated remediation action triggered instantly", GREEN),
    ("📱 Notify", "Multilingual SMS to affected customers (8 languages)", ACCENT),
]
box_w = Inches(2.3)
box_h = Inches(1.5)
arrow_w = Inches(0.3)
total = len(steps)
start_x = Inches(0.4)

for i, (icon_title, desc, col) in enumerate(steps):
    x = start_x + i * (box_w + arrow_w)
    add_rect(s, x, Inches(2.0), box_w, box_h, DARK_CARD)
    add_rect(s, x, Inches(2.0), box_w, Inches(0.05), col)
    add_textbox(s, icon_title, x, Inches(2.1), box_w, Inches(0.4),
                font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + Inches(0.1), Inches(2.55), box_w - Inches(0.2), Inches(0.85),
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow
    if i < total - 1:
        ax = x + box_w + Inches(0.05)
        add_textbox(s, "▶", ax, Inches(2.55), arrow_w, Inches(0.4),
                    font_size=16, color=GREY, align=PP_ALIGN.CENTER)

# Key differentiators
diffs = [
    ("< 2 seconds", "End-to-end pipeline latency"),
    ("8 languages", "Customer notifications"),
    ("0 human steps", "For remediable faults"),
    ("3 AI services", "Classify · Predict · Detect"),
]
for i, (val, label) in enumerate(diffs):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(4.0), Inches(3.0), Inches(1.1), DARK_CARD)
    add_textbox(s, val, x, Inches(4.05), Inches(3.0), Inches(0.5),
                font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, Inches(4.55), Inches(3.0), Inches(0.4),
                font_size=11, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(s, "PULSE handles the entire lifecycle automatically — from first log to customer notification.",
            Inches(0.5), Inches(5.3), Inches(12), Inches(0.4),
            font_size=13, color=GREY, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "System Architecture", "Three-service microservice architecture")

# Three service boxes
services = [
    ("React Frontend :3000", ["React 19 + TypeScript", "Redux Toolkit / RTK Query", "Leaflet maps", "Vite build", "WebSocket client"], ACCENT),
    ("Django Backend :8000", ["Django 6 REST API", "Django Channels (WS)", "SimpleJWT auth", "SQLite database", "Background threads"], ACCENT2),
    ("FastAPI AI :8001", ["/classify — root cause", "/predict — failure prob.", "/detect — Z-score anomaly", "Response < 50ms", "Stateless microservice"], GREEN),
]
for i, (title, items, col) in enumerate(services):
    x = Inches(0.4 + i * 4.3)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(3.2), DARK_CARD)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(0.06), col)
    add_textbox(s, title, x + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.45),
                font_size=14, bold=True, color=col)
    for j, item in enumerate(items):
        add_textbox(s, f"• {item}",
                    x + Inches(0.15), Inches(2.0) + j * Inches(0.42),
                    Inches(3.7), Inches(0.38),
                    font_size=11, color=WHITE)
    # Arrow between services
    if i < 2:
        add_textbox(s, "HTTP\n⟶",
                    x + Inches(4.0), Inches(2.6),
                    Inches(0.3), Inches(0.5),
                    font_size=9, color=GREY, align=PP_ALIGN.CENTER)

# Data flow below
add_textbox(s, "Data Flow:",
            Inches(0.4), Inches(4.85), Inches(2), Inches(0.35),
            font_size=13, bold=True, color=ACCENT2)
add_textbox(s,
    "ATM/Channel → POST /api/logs/ingest/ → LogEntry saved → process_log() thread → "
    "[AI classify] → [Z-score detect] → Incident → Health Score → Alert → Self-Heal → Customer SMS → WebSocket broadcast",
    Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.0),
    font_size=11, color=WHITE)

add_textbox(s, "No Celery / Redis required — lightweight daemon threads handle all async processing",
            Inches(0.4), Inches(6.3), Inches(12), Inches(0.4),
            font_size=11, color=GREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Tech Stack", "Modern, production-ready technologies")

categories = [
    ("Frontend", [
        "React 19 — latest concurrent features",
        "TypeScript — end-to-end type safety",
        "Redux Toolkit + RTK Query — state & API layer",
        "React-Leaflet — interactive India ATM map",
        "Vite — sub-second HMR dev server",
        "Tailwind CSS — utility-first dark UI",
    ], ACCENT),
    ("Backend", [
        "Django 6 — batteries-included REST framework",
        "Django REST Framework — serialization & auth",
        "Django Channels — WebSocket support (ASGI)",
        "SimpleJWT — stateless JWT authentication",
        "SQLite — zero-config embedded database",
        "threading.Thread — lightweight async processing",
    ], ACCENT2),
    ("AI / Data", [
        "FastAPI + Uvicorn — high-performance AI microservice",
        "Rule-based classifier — 95%+ accuracy on known faults",
        "Z-score anomaly detection — no ML training required",
        "Rolling window prediction — failure probability scoring",
        "Pydantic v2 — request/response validation",
        "Python 3.14 — latest CPython runtime",
    ], GREEN),
]

for i, (cat, items, col) in enumerate(categories):
    x = Inches(0.4 + i * 4.3)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(4.5), DARK_CARD)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(0.06), col)
    add_textbox(s, cat, x + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.45),
                font_size=16, bold=True, color=col)
    for j, item in enumerate(items):
        add_textbox(s, f"• {item}",
                    x + Inches(0.15), Inches(2.05) + j * Inches(0.55),
                    Inches(3.7), Inches(0.5),
                    font_size=10.5, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI Pipeline Deep Dive
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "AI Pipeline — How It Works", "Three AI modules, one unified intelligence layer")

modules = [
    ("/classify", "Root Cause Classifier",
     ["Input: eventCode + logLevel",
      "Maps to 8 root cause categories",
      "Returns category + confidence + detail",
      "NETWORK · SWITCH · SERVER · TIMEOUT",
      "CASH_JAM · FRAUD · HARDWARE · UNKNOWN",
      "Response time: < 10ms"],
     ACCENT,
     "Incident only created\nif confidence ≥ 65%"),
    ("/detect", "Z-Score Anomaly Detector",
     ["Input: current 1h error rate",
      "vs historical baseline mean/std",
      "Z-score threshold: 2.0σ",
      "Flags: RAPID_FAILURES",
      "UNUSUAL_WITHDRAWAL etc.",
      "Demo fallback: 5% baseline"],
     ACCENT2,
     "Works from first\nminute of data"),
    ("/predict", "Failure Probability Predictor",
     ["Input: health snapshot time-series",
      "Rolling mean + variance + slope",
      "Outputs failure probability 0–1",
      "Used in AI Analysis dashboard",
      "Per-ATM future risk score",
      "No training data needed"],
     GREEN,
     "Proactive, not\njust reactive"),
]

for i, (endpoint, title, bullets, col, note) in enumerate(modules):
    x = Inches(0.4 + i * 4.3)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(4.2), DARK_CARD)
    add_rect(s, x, Inches(1.4), Inches(4.0), Inches(0.06), col)
    add_textbox(s, endpoint, x + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.35),
                font_size=12, bold=True, color=col)
    add_textbox(s, title, x + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.4),
                font_size=13, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        add_textbox(s, f"• {b}",
                    x + Inches(0.15), Inches(2.3) + j * Inches(0.42),
                    Inches(3.7), Inches(0.38),
                    font_size=10, color=WHITE)
    # Note badge
    add_rect(s, x + Inches(0.15), Inches(5.1), Inches(3.7), Inches(0.55), col)
    add_textbox(s, note, x + Inches(0.15), Inches(5.1), Inches(3.7), Inches(0.55),
                font_size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Self-Healing Engine
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Self-Healing Engine", "Automated remediation — no human in the loop")

# Self-heal map table
add_rect(s, Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5), DARK_CARD)
add_rect(s, Inches(0.5), Inches(1.4), Inches(8.5), Inches(0.06), ACCENT)

headers = ["Root Cause", "Action", "Resolution", "Human Needed?"]
col_widths = [Inches(2.0), Inches(2.5), Inches(2.5), Inches(1.5)]
col_x = [Inches(0.6), Inches(2.6), Inches(5.1), Inches(7.6)]

for j, (hdr, cx) in enumerate(zip(headers, col_x)):
    add_textbox(s, hdr, cx, Inches(1.5), col_widths[j], Inches(0.4),
                font_size=12, bold=True, color=ACCENT2)

rows = [
    ("NETWORK",  "SWITCH_NETWORK",   "AUTO_RESOLVED",   "❌ No",  GREEN),
    ("SWITCH",   "RESTART_SERVICE",  "AUTO_RESOLVED",   "❌ No",  GREEN),
    ("SERVER",   "RESTART_SERVICE",  "AUTO_RESOLVED",   "❌ No",  GREEN),
    ("TIMEOUT",  "FLUSH_CACHE",      "AUTO_RESOLVED",   "❌ No",  GREEN),
    ("FRAUD",    "FREEZE_ATM",       "Manual Review",   "✅ Yes", AMBER),
    ("CASH_JAM", "ALERT_ENGINEER",   "Field Visit",     "✅ Yes", AMBER),
    ("HARDWARE", "ALERT_ENGINEER",   "Field Visit",     "✅ Yes", AMBER),
]

for k, (cause, action, resolution, human, col) in enumerate(rows):
    y = Inches(2.0) + k * Inches(0.48)
    row_bg = RGBColor(0x12, 0x12, 0x1e) if k % 2 == 0 else DARK_CARD
    add_rect(s, Inches(0.5), y, Inches(8.5), Inches(0.46), row_bg)
    vals = [cause, action, resolution, human]
    for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_widths)):
        color = col if j == 0 else (GREEN if human == "❌ No" and j == 3 else WHITE)
        if j == 3:
            color = GREEN if "No" in human else AMBER
        add_textbox(s, val, cx, y + Inches(0.05), cw, Inches(0.36),
                    font_size=11, color=color)

# Right side callout
add_rect(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(4.5), DARK_CARD)
add_rect(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(0.06), GREEN)
add_textbox(s, "Impact", Inches(9.4), Inches(1.5), Inches(3.5), Inches(0.4),
            font_size=14, bold=True, color=GREEN)

callouts = [
    ("57%", "of faults auto-resolved"),
    ("< 2s", "remediation time"),
    ("0 engineers", "for network/software faults"),
    ("24/7", "autonomous operation"),
]
for i, (val, label) in enumerate(callouts):
    y = Inches(2.0) + i * Inches(0.88)
    add_textbox(s, val, Inches(9.4), y, Inches(3.5), Inches(0.45),
                font_size=22, bold=True, color=GREEN)
    add_textbox(s, label, Inches(9.4), y + Inches(0.42), Inches(3.5), Inches(0.35),
                font_size=11, color=GREY)

add_textbox(s, "Auto-resolved incidents have status AUTO_RESOLVED set in < 2 seconds",
            Inches(0.5), Inches(6.15), Inches(12), Inches(0.4),
            font_size=11, color=GREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Anomaly Detection
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Anomaly Detection", "Z-score statistical analysis on real-time error rates")

# Formula display
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.1), DARK_CARD)
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.05), ACCENT2)
add_textbox(s, "Algorithm:",
            Inches(0.7), Inches(1.48), Inches(3), Inches(0.35),
            font_size=13, bold=True, color=ACCENT2)
add_textbox(s,
    "Z-score = (current_error_rate − historical_mean) / historical_std_dev     →    If Z > 2.0 : FLAG as ANOMALY",
    Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
    font_size=14, bold=True, color=WHITE)

# Two columns
add_card(s, Inches(0.5), Inches(2.7), Inches(6.0), Inches(3.0),
         "Detection Process", [
             "Compute current 1-hour error rate",
             "Compare against all historical logs (pre-window)",
             "If < 5 baseline logs: use 5% mean / 1.5% std (demo default)",
             "Anomaly flagged if Z-score ≥ 2.0",
             "AnomalyFlag record created with type + confidence",
         ], ACCENT2)

add_card(s, Inches(6.8), Inches(2.7), Inches(6.0), Inches(3.0),
         "Anomaly Types Detected", [
             "RAPID_FAILURES — error rate spike",
             "UNUSUAL_WITHDRAWAL — transaction pattern anomaly",
             "CARD_SKIMMING — suspicious card read sequences",
             "MALWARE_PATTERN — unusual event code combinations",
             "All flagged with confidence score 0.0 – 1.0",
         ], AMBER)

add_textbox(s, "Works from first minute — no historical data required for demo mode",
            Inches(0.5), Inches(6.15), Inches(12), Inches(0.4),
            font_size=12, color=GREEN, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Health Score System
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Health Score System", "Real-time ATM health computed from live log stream")

# Formula
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.3), DARK_CARD)
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.05), GREEN)
add_textbox(s, "Health Score Algorithm (§11):",
            Inches(0.7), Inches(1.48), Inches(5), Inches(0.35),
            font_size=13, bold=True, color=GREEN)
add_textbox(s, "score = 100  −  min(40, error_rate × 100)  −  min(30, warns × 2)  −  (open_incidents × 15)",
            Inches(0.7), Inches(1.85), Inches(12), Inches(0.45),
            font_size=14, bold=True, color=WHITE)
add_textbox(s, "where error_rate = (errors + criticals × 2) / total_logs_in_last_1_hour",
            Inches(0.7), Inches(2.3), Inches(12), Inches(0.35),
            font_size=11, color=GREY)

# Sub-scores
sub_scores = [
    ("Network Score", "Connectivity & routing health", ACCENT),
    ("Hardware Score", "Physical components health", RED),
    ("Software Score", "OS + application stability", ACCENT2),
    ("Transaction Score", "Payment success rate", GREEN),
]
for i, (title, desc, col) in enumerate(sub_scores):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(3.0), Inches(3.0), Inches(1.1), DARK_CARD)
    add_rect(s, x, Inches(3.0), Inches(3.0), Inches(0.05), col)
    add_textbox(s, title, x + Inches(0.1), Inches(3.1), Inches(2.8), Inches(0.4),
                font_size=12, bold=True, color=col)
    add_textbox(s, desc, x + Inches(0.1), Inches(3.55), Inches(2.8), Inches(0.4),
                font_size=10, color=GREY)
add_textbox(s, "Overall health = mean of 4 sub-scores. Each INFO log +2 pts recovery. Each CRITICAL -10 pts.",
            Inches(0.5), Inches(4.25), Inches(12.3), Inches(0.4),
            font_size=11, color=WHITE)

# Status thresholds
thresholds = [
    ("≥ 75", "ONLINE", GREEN),
    ("30–74", "DEGRADED", AMBER),
    ("< 30", "OFFLINE", RED),
]
add_textbox(s, "ATM Status Thresholds:", Inches(0.5), Inches(4.9), Inches(4), Inches(0.4),
            font_size=13, bold=True, color=ACCENT2)
for i, (range_, status, col) in enumerate(thresholds):
    add_rect(s, Inches(0.5 + i * 4.1), Inches(5.4), Inches(3.8), Inches(0.8), col)
    add_textbox(s, f"Score {range_}  →  {status}",
                Inches(0.5 + i * 4.1), Inches(5.4), Inches(3.8), Inches(0.8),
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Live Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Live Dashboard", "Real-time fleet intelligence at a glance")

features = [
    ("KPI Cards", [
        "Total ATMs in fleet",
        "Online / Degraded / Offline counts",
        "Open incidents count",
        "Average health score",
        "Auto-refresh every 30s",
    ], Inches(0.4), Inches(1.4), ACCENT),
    ("ATM Network Map", [
        "India-only Leaflet map",
        "Color-coded markers by health",
        "Click → side panel details",
        "Click → full ATM detail page",
        "Real-time marker updates",
    ], Inches(4.6), Inches(1.4), ACCENT2),
    ("Health Overview", [
        "Fleet-wide health distribution",
        "Online / Degraded / Offline bars",
        "Top 5 worst ATMs list",
        "Health trend over time",
        "Polled every 30 seconds",
    ], Inches(8.8), Inches(1.4), GREEN),
    ("Incident Feed", [
        "Live recent incidents",
        "Severity color-coded badges",
        "One-click resolve button",
        "AI root cause displayed",
        "Self-heal action shown",
    ], Inches(0.4), Inches(4.2), AMBER),
    ("Self-Heal Log", [
        "Recent auto-remediation actions",
        "Action type + outcome",
        "Linked to source incident",
        "Timestamp + ATM name",
        "Filter by action type",
    ], Inches(4.6), Inches(4.2), ACCENT),
    ("AI Pipeline Feed", [
        "Last 40 processed log events",
        "Classification + confidence",
        "Incident created indicator",
        "Self-heal action column",
        "REST poll every 3 seconds",
    ], Inches(8.8), Inches(4.2), ACCENT2),
]

for title, bullets, x, y, col in features:
    add_rect(s, x, y, Inches(4.0), Inches(2.5), DARK_CARD)
    add_rect(s, x, y, Inches(4.0), Inches(0.05), col)
    add_textbox(s, title, x + Inches(0.1), y + Inches(0.1), Inches(3.8), Inches(0.4),
                font_size=13, bold=True, color=col)
    for j, b in enumerate(bullets):
        add_textbox(s, f"• {b}", x + Inches(0.1), y + Inches(0.55) + j * Inches(0.38),
                    Inches(3.8), Inches(0.35), font_size=10, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Multilingual Notifications
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Multilingual Customer Notifications", "India's 8 major languages, automated")

# Languages grid
langs = ["English", "हिन्दी (Hindi)", "தமிழ் (Tamil)", "తెలుగు (Telugu)",
         "ಕನ್ನಡ (Kannada)", "मराठी (Marathi)", "বাংলা (Bengali)", "ગુજરાતી (Gujarati)"]
colors_l = [ACCENT, ACCENT2, GREEN, AMBER, RED, ACCENT, ACCENT2, GREEN]

for i, (lang, col) in enumerate(zip(langs, colors_l)):
    row, c = divmod(i, 4)
    x = Inches(0.4 + c * 3.2)
    y = Inches(1.4 + row * 0.85)
    add_rect(s, x, y, Inches(3.0), Inches(0.65), DARK_CARD)
    add_rect(s, x, y, Inches(0.06), Inches(0.65), col)
    add_textbox(s, lang, x + Inches(0.15), y + Inches(0.1), Inches(2.8), Inches(0.45),
                font_size=13, bold=True, color=WHITE)

# Template system
add_rect(s, Inches(0.4), Inches(3.4), Inches(12.5), Inches(2.0), DARK_CARD)
add_rect(s, Inches(0.4), Inches(3.4), Inches(12.5), Inches(0.06), ACCENT)
add_textbox(s, "40 Message Templates (8 languages × 5 incident types):",
            Inches(0.6), Inches(3.5), Inches(11), Inches(0.4),
            font_size=13, bold=True, color=ACCENT2)
template_types = ["atm_offline", "upi_timeout", "card_decline", "network_failure", "cash_jam"]
for i, t in enumerate(template_types):
    add_textbox(s, f"• {t}",
                Inches(0.6 + i * 2.4), Inches(3.95), Inches(2.3), Inches(0.35),
                font_size=11, color=WHITE)
add_textbox(s, "Templates seeded via: python manage.py seed_templates",
            Inches(0.6), Inches(4.45), Inches(11), Inches(0.35),
            font_size=10, color=GREY, italic=True)

# Example SMS
add_rect(s, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.2), RGBColor(0x06, 0x10, 0x08))
add_rect(s, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.05), GREEN)
add_textbox(s, "📱 Example (English): Dear customer, ATM INC-3F8A1B2C at Chennai Main Branch is temporarily offline. "
               "Our team has been notified and is working to restore service. We apologize for the inconvenience.",
            Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.9),
            font_size=11, color=WHITE, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Data Models & API
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Data Models & API Coverage", "31 REST endpoints, 10 data models, 100% connected")

# Models
models_list = [
    ("ATM", "id, name, location, status, healthScore, networkScore, hardwareScore, softwareScore, transactionScore, lat/lon", ACCENT),
    ("LogEntry", "sourceId (UUID), sourceType, logLevel, eventCode, message, processed, timestamp", ACCENT2),
    ("Incident", "incidentId, title, severity, status, rootCauseCategory, aiConfidence, triggeringLogId", GREEN),
    ("HealthSnapshot", "sourceId, healthScore, status, sub-scores, timestamp", AMBER),
    ("Alert", "incidentId, alertType, title, severity, acknowledged, sentAt", RED),
    ("SelfHealAction", "incidentId, actionType, triggeredBy, status, result", GREEN),
    ("AnomalyFlag", "sourceId, anomalyType, confidenceScore, status, description", AMBER),
    ("MessageTemplate", "templateKey, language, channel, body — unique(key+lang+channel)", ACCENT),
    ("CustomerNotification", "recipientId, channel, message, language, status, sentAt", ACCENT2),
]
add_textbox(s, "Core Data Models:", Inches(0.5), Inches(1.4), Inches(4), Inches(0.4),
            font_size=14, bold=True, color=ACCENT2)
for i, (model, fields, col) in enumerate(models_list):
    y = Inches(1.85) + i * Inches(0.54)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5), DARK_CARD)
    add_textbox(s, model, Inches(0.6), y + Inches(0.05), Inches(2.0), Inches(0.4),
                font_size=11, bold=True, color=col)
    add_textbox(s, fields, Inches(2.65), y + Inches(0.05), Inches(10.0), Inches(0.4),
                font_size=9.5, color=GREY)

add_textbox(s, "31 API endpoints · 31 RTK Query hooks · 100% matched (no orphaned endpoints)",
            Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
            font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Results & Impact
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Results & Impact", "What PULSE delivers in production")

# Big metrics
metrics = [
    ("< 2 sec", "End-to-end pipeline\nlatency (ingest → notify)", GREEN),
    ("57%", "Faults auto-resolved\nwithout human intervention", ACCENT),
    ("8", "Languages supported\nfor customer SMS", ACCENT2),
    ("3 AI", "Services: classify +\npredict + detect", AMBER),
    ("40", "Message templates\nseeded across languages", ACCENT),
    ("0", "Redis/Celery dependencies\n(pure Python threading)", GREEN),
]
for i, (val, label, col) in enumerate(metrics):
    row, c = divmod(i, 3)
    x = Inches(0.5 + c * 4.25)
    y = Inches(1.5 + row * 2.2)
    add_rect(s, x, y, Inches(4.0), Inches(1.9), DARK_CARD)
    add_rect(s, x, y, Inches(4.0), Inches(0.06), col)
    add_textbox(s, val, x, y + Inches(0.15), Inches(4.0), Inches(0.75),
                font_size=42, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, y + Inches(0.95), Inches(4.0), Inches(0.75),
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "Fully functional end-to-end demo: ATM → Log → AI → Incident → Self-Heal → Customer SMS",
            Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
            font_size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion & Future Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_heading(s, "Conclusion & Roadmap", "PULSE — Redefining ATM intelligence for India")

# Left: summary
add_rect(s, Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.0), DARK_CARD)
add_rect(s, Inches(0.4), Inches(1.4), Inches(6.0), Inches(0.06), ACCENT)
add_textbox(s, "What We Built", Inches(0.6), Inches(1.5), Inches(5.7), Inches(0.4),
            font_size=14, bold=True, color=ACCENT)
summary_points = [
    "✅ End-to-end ATM monitoring platform",
    "✅ AI root cause classifier (3 services)",
    "✅ Autonomous self-healing engine",
    "✅ Z-score anomaly detection",
    "✅ Health score with 4 dimensions",
    "✅ Real-time React dashboard + Leaflet map",
    "✅ 40 multilingual notification templates",
    "✅ JWT auth + RTK Query + WebSocket",
    "✅ Synthetic data simulator for demo",
    "✅ 31 REST API endpoints, 100% connected",
]
for i, pt in enumerate(summary_points):
    add_textbox(s, pt, Inches(0.6), Inches(2.0) + i * Inches(0.45),
                Inches(5.7), Inches(0.4), font_size=11, color=WHITE)

# Right: roadmap
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.0), DARK_CARD)
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.06), GREEN)
add_textbox(s, "Future Roadmap", Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
            font_size=14, bold=True, color=GREEN)
roadmap = [
    ("v2 — ML Model", "Train on real ATM fault datasets (LSTM/XGBoost)"),
    ("v2 — Redis Pub/Sub", "Replace InMemoryChannelLayer for true WebSocket"),
    ("v2 — Mobile App", "React Native app for field engineers"),
    ("v3 — Edge AI", "Deploy classifier on-device at ATM controller"),
    ("v3 — Predictive", "72-hour failure prediction from sensor telemetry"),
    ("v3 — Multi-bank", "SaaS platform — serve 100+ banks on single infra"),
    ("v3 — WhatsApp", "WhatsApp Business API for customer notifications"),
]
for i, (phase, desc) in enumerate(roadmap):
    y = Inches(2.0) + i * Inches(0.55)
    add_textbox(s, phase, Inches(7.0), y, Inches(1.8), Inches(0.45),
                font_size=11, bold=True, color=GREEN)
    add_textbox(s, desc, Inches(8.9), y, Inches(3.7), Inches(0.45),
                font_size=10, color=WHITE)

# Bottom tagline
add_rect(s, 0, H - Inches(0.7), W, Inches(0.7), ACCENT)
add_textbox(s, "PULSE — Detect · Classify · Heal · Notify · Repeat",
            0, H - Inches(0.65), W, Inches(0.6),
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\Praanesh MB\OneDrive\Documents\PULSE\PULSE_Presentation.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
